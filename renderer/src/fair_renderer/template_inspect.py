"""fair-template: bind an arbitrary PowerPoint template to the FAIR profile.

The layout map is the adapter between region names sessions author
against and whatever a designer's template actually contains
(`layoutmap.py`). Writing that map by hand is the one step that stops
"bring your own template" from working, because a real template's
placeholder indices are not guessable.

This inspects a .pptx, matches its layouts against the Core profile
(docs/template-profile.md), infers each region, and emits a draft
layout-map.yaml plus a conformance report.

    fair-template their-template.pptx --out layout-map.yaml

Inference reads placeholder TYPE and GEOMETRY, never index order. A
designer can legitimately ship Two Content with the right-hand
placeholder at idx 1; assuming 1 == left would silently mirror every
split slide in the corpus, and nothing downstream would notice.
"""

from __future__ import annotations

import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

# Never mapped as regions: PowerPoint's own slide furniture, present on
# nearly every layout at these fixed indices.
CHROME = {
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
}
TITLES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


class TemplateError(Exception):
    pass


@dataclass
class Placeholder:
    idx: int
    type: object
    left: int | None
    top: int | None
    width: int | None
    height: int | None

    @property
    def area(self) -> int:
        if self.width is None or self.height is None:
            return 0
        return self.width * self.height

    @property
    def positioned(self) -> bool:
        return self.left is not None and self.top is not None


@dataclass
class LayoutResult:
    key: str
    layout_name: str | None = None
    regions: dict[str, int] = field(default_factory=dict)
    missing: list[str] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return self.layout_name is not None and not self.missing


def _placeholders(layout) -> list[Placeholder]:
    out = []
    for ph in layout.placeholders:
        fmt = ph.placeholder_format
        out.append(
            Placeholder(
                idx=fmt.idx,
                type=fmt.type,
                left=ph.left,
                top=ph.top,
                width=ph.width,
                height=ph.height,
            )
        )
    return out


def _split(phs: list[Placeholder]) -> tuple[Placeholder | None, list[Placeholder]]:
    """(title, content) with chrome discarded."""
    title = next((p for p in phs if p.type in TITLES), None)
    content = [p for p in phs if p.type not in TITLES and p.type not in CHROME]
    # Stable order even when a layout leaves placeholders unpositioned.
    content.sort(key=lambda p: (p.top or 0, p.left or 0, p.idx))
    return title, content


def _by_x(phs: list[Placeholder]) -> list[Placeholder]:
    """Left-to-right, falling back to index when geometry is inherited."""
    if all(p.positioned for p in phs):
        return sorted(phs, key=lambda p: p.left)
    return sorted(phs, key=lambda p: p.idx)


# --- per-layout resolvers -------------------------------------------------
#
# Each takes (title, content) and returns {region: idx}. A region it
# cannot bind is simply absent; the caller reports it as missing.


def _r_title_subtitle(title, content):
    out = {}
    if title:
        out["title"] = title.idx
    sub = next((p for p in content if p.type == PP_PLACEHOLDER.SUBTITLE), None)
    sub = sub or (content[0] if content else None)
    if sub:
        out["subtitle"] = sub.idx
    return out


def _r_single(region: str):
    def resolve(title, content):
        out = {}
        if title:
            out["title"] = title.idx
        if content:
            # The largest body is the content well; designers often leave
            # a small decorative placeholder alongside it.
            out[region] = max(content, key=lambda p: p.area).idx if any(
                p.area for p in content
            ) else content[0].idx
        return out

    return resolve


def _r_two_content(title, content):
    out = {}
    if title:
        out["title"] = title.idx
    if len(content) >= 2:
        left, right = _by_x(content)[:2]
        out["left"], out["right"] = left.idx, right.idx
    return out


def _r_comparison(title, content):
    out = {}
    if title:
        out["title"] = title.idx
    if len(content) < 4:
        return out
    ordered = _by_x(content)
    left_side, right_side = ordered[:2], ordered[2:4]
    for side, prefix in ((left_side, "left"), (right_side, "right")):
        # The heading sits above its body; without geometry, lower idx.
        if all(p.positioned for p in side):
            head, body = sorted(side, key=lambda p: p.top)
        else:
            head, body = sorted(side, key=lambda p: p.idx)
        out[f"{prefix}_head"] = head.idx
        out[prefix] = body.idx
    return out


def _r_caption(title, content):
    out = {}
    if title:
        out["title"] = title.idx
    if len(content) >= 2:
        # The content well is the bigger of the two; the caption is prose.
        ranked = sorted(content, key=lambda p: p.area, reverse=True)
        out["content"], out["caption"] = ranked[0].idx, ranked[1].idx
    return out


def _r_picture(title, content):
    out = {}
    if title:
        out["title"] = title.idx
    pic = next((p for p in content if p.type == PP_PLACEHOLDER.PICTURE), None)
    if pic:
        out["picture"] = pic.idx
    rest = [p for p in content if pic is None or p.idx != pic.idx]
    if rest:
        out["caption"] = max(rest, key=lambda p: p.area).idx
    return out


def _r_title_only(title, content):
    return {"title": title.idx} if title else {}


def _r_none(title, content):
    return {}


@dataclass(frozen=True)
class Spec:
    key: str
    aliases: tuple[str, ...]
    regions: tuple[str, ...]
    resolve: object


# The Core profile. Every one of these exists in any template built on
# PowerPoint's default masters, under one of the listed names.
CORE: tuple[Spec, ...] = (
    Spec("Title", ("Title Slide", "Title"), ("title", "subtitle"), _r_title_subtitle),
    Spec("Full", ("Title and Content", "Full", "Content"), ("title", "full"), _r_single("full")),
    Spec("Section", ("Section Header", "Section"), ("title", "subtitle"), _r_title_subtitle),
    Spec("Split", ("Two Content", "Split", "Comparison "), ("title", "left", "right"), _r_two_content),
    Spec(
        "Comparison",
        ("Comparison",),
        ("title", "left_head", "left", "right_head", "right"),
        _r_comparison,
    ),
    Spec(
        "Caption",
        ("Content with Caption", "Caption"),
        ("title", "content", "caption"),
        _r_caption,
    ),
    Spec(
        "Picture",
        ("Picture with Caption", "Picture"),
        ("title", "picture", "caption"),
        _r_picture,
    ),
    Spec("TitleOnly", ("Title Only", "TitleOnly"), ("title",), _r_title_only),
    Spec("Blank", ("Blank",), (), _r_none),
)

# Optional. Absent is fine; a session using one fails validation loudly.
EXTENDED: tuple[Spec, ...] = (
    Spec(
        "Cards",
        ("Cards",),
        ("title", "head1", "head2", "head3", "head4", "card1", "card2", "card3", "card4"),
        None,  # bespoke: bound by declared order, see _resolve_cards
    ),
)


def _resolve_cards(title, content):
    out = {}
    if title:
        out["title"] = title.idx
    if len(content) < 8:
        return out
    columns = _by_x(content)
    # Four columns of two: the tab sits above its body panel.
    groups = [columns[i : i + 2] for i in range(0, 8, 2)]
    for n, group in enumerate(groups, start=1):
        if all(p.positioned for p in group):
            head, body = sorted(group, key=lambda p: p.top)
        else:
            head, body = sorted(group, key=lambda p: p.idx)
        out[f"head{n}"] = head.idx
        out[f"card{n}"] = body.idx
    return out


def _match(layouts: dict[str, object], spec: Spec):
    for alias in spec.aliases:
        for name, layout in layouts.items():
            if name.strip().lower() == alias.strip().lower():
                return name, layout
    return None, None


def _columns(phs: list[Placeholder], count: int) -> list[list[Placeholder]] | None:
    """Group placeholders into `count` columns by x, or None if they don't."""
    if not all(p.positioned for p in phs):
        return None
    groups: list[list[Placeholder]] = []
    for ph in sorted(phs, key=lambda p: p.left):
        # Same column when the left edges are within a placeholder's width.
        if groups and abs(ph.left - groups[-1][0].left) < (ph.width or 0) * 0.5:
            groups[-1].append(ph)
        else:
            groups.append([ph])
    return groups if len(groups) == count else None


def _find_cards(layouts: dict[str, object], claimed: set[str]):
    """Find a four-panel layout under whatever name the designer gave it.

    A template that already has a quad-panel layout should be used, not
    made to grow a second one called "Cards". Match on shape: four
    columns of two content placeholders, a heading above each panel.
    """
    best = None
    for name, layout in layouts.items():
        if name in claimed:
            continue
        title, content = _split(_placeholders(layout))
        if len(content) < 8:
            continue
        groups = _columns(content, 4)
        if not groups or any(len(g) < 2 for g in groups):
            continue
        bound = _resolve_cards(title, content)
        if len({k for k in bound if k.startswith(("head", "card"))}) == 8:
            # Prefer the tightest fit: exactly eight content placeholders.
            score = abs(len(content) - 8)
            if best is None or score < best[0]:
                best = (score, name, layout, bound)
    if best is None:
        return None, None, None
    _, name, layout, bound = best
    return name, layout, bound


def inspect_template(path: Path) -> list[LayoutResult]:
    prs = Presentation(str(path))
    layouts = {
        layout.name: layout
        for master in prs.slide_masters
        for layout in master.slide_layouts
    }
    results: list[LayoutResult] = []
    claimed: set[str] = set()
    for spec in CORE + EXTENDED:
        result = LayoutResult(key=spec.key)
        name, layout = _match(layouts, spec)
        if layout is None and spec.key == "Cards":
            # Fall back to shape: the four-panel layout may be called
            # anything, and using theirs beats requiring ours.
            name, layout, bound = _find_cards(layouts, claimed)
            if layout is not None:
                result.layout_name = name
                result.regions = bound
                result.missing = [r for r in spec.regions if r not in bound]
                result.notes.append(f"matched by shape, not name: {name!r}")
                claimed.add(name)
                results.append(result)
                continue
        if layout is None:
            result.missing = list(spec.regions)
            result.notes.append(
                f"no layout named {' or '.join(repr(a) for a in spec.aliases)}"
            )
            results.append(result)
            continue
        claimed.add(name)
        result.layout_name = name
        title, content = _split(_placeholders(layout))
        resolve = _resolve_cards if spec.key == "Cards" else spec.resolve
        result.regions = resolve(title, content)
        result.missing = [r for r in spec.regions if r not in result.regions]
        if spec.key == "Picture" and "picture" not in result.regions:
            result.notes.append(
                "no PICTURE placeholder: images here would be free-positioned, "
                "not template-bound"
            )
        results.append(result)
    return results


def to_layout_map(results: list[LayoutResult]) -> str:
    """Render the bindings as layout-map.yaml, in profile order."""
    lines = ["# Generated by fair-template. Region -> placeholder index.",
             "# Regenerate after editing the template; the build validates it."]
    for result in results:
        if result.layout_name is None:
            continue
        lines.append(f"{result.key}:")
        lines.append(f'  layout_name: "{result.layout_name}"')
        if not result.regions:
            lines.append("  regions: {}")
            continue
        lines.append("  regions:")
        for region in result.regions:
            lines.append(f"    {region}: {result.regions[region]}")
    return "\n".join(lines) + "\n"


def unmapped(path: Path, results: list[LayoutResult]) -> list[tuple[str, int, list[str]]]:
    """Layouts the profile did not claim, with what they contain.

    A designer's template usually carries more than the profile needs.
    Listing them means a library can map its own extra keys by hand
    instead of discovering them by accident.
    """
    prs = Presentation(str(path))
    used = {r.layout_name for r in results if r.layout_name}
    out = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name in used:
                continue
            title, content = _split(_placeholders(layout))
            kinds = [str(p.type).split()[0].lower() for p in content]
            out.append((layout.name, len(content), kinds))
    return out


def report(results: list[LayoutResult]) -> str:
    core_keys = {s.key for s in CORE}
    lines, problems = [], 0
    for result in results:
        tier = "core" if result.key in core_keys else "extended"
        if result.layout_name is None:
            mark = "MISSING" if tier == "core" else "absent"
            if tier == "core":
                problems += 1
            lines.append(f"  {mark:8} {result.key:12} ({tier}) — {result.notes[0]}")
            continue
        if result.missing:
            if tier == "core":
                problems += 1
            lines.append(
                f"  PARTIAL  {result.key:12} -> {result.layout_name!r}; "
                f"unbound: {', '.join(result.missing)}"
            )
        else:
            bound = ", ".join(f"{k}={v}" for k, v in result.regions.items())
            lines.append(f"  ok       {result.key:12} -> {result.layout_name!r}  {bound}")
        for note in result.notes:
            lines.append(f"           note: {note}")
    head = (
        "Template conforms to the Core profile."
        if problems == 0
        else f"{problems} core layout(s) unusable — see docs/template-profile.md"
    )
    return head + "\n" + "\n".join(lines)


def main(argv: list[str] | None = None) -> int:
    import argparse

    ap = argparse.ArgumentParser(
        prog="fair-template",
        description="Inspect a .pptx and emit a draft layout-map.yaml plus a conformance report",
    )
    ap.add_argument("template", type=Path)
    ap.add_argument(
        "--add-cards",
        type=Path,
        metavar="OUT",
        default=None,
        help="write a copy of the template with a Cards layout added, in its own theme colours",
    )
    ap.add_argument("--force", action="store_true",
                    help="add Cards even if the template already has a four-panel layout")
    ap.add_argument("--out", type=Path, default=None,
                    help="write layout-map.yaml here (default: stdout)")
    ap.add_argument("--quiet", action="store_true", help="suppress the report")
    args = ap.parse_args(argv)

    if not args.template.exists():
        print(f"error: no such template: {args.template}", file=sys.stderr)
        return 1

    if args.add_cards:
        from .cards_layout import CardsError, inject_cards

        try:
            outcome = inject_cards(args.template, args.add_cards, force=args.force)
        except CardsError as e:
            print(f"error: {e}", file=sys.stderr)
            return 1
        if outcome["added"]:
            print(
                f"added the Cards layout to {args.add_cards} "
                f"(cloned from {outcome['donor']!r}, filled from the template's own theme)"
            )
        else:
            print(
                f"{args.template} already has a four-panel layout "
                f"({outcome['layout']!r}); copied unchanged to {args.add_cards}. "
                "Use --force to add another."
            )
        # Report against the file we just wrote, not the original.
        args.template = args.add_cards

    results = inspect_template(args.template)

    if args.out:
        args.out.write_text(to_layout_map(results), encoding="utf-8")
        if not args.quiet:
            print(f"wrote {args.out}")
    else:
        print(to_layout_map(results))

    if not args.quiet:
        print(report(results), file=sys.stderr)
        extra = unmapped(args.template, results)
        if extra:
            print("\nAlso in this template, unmapped:", file=sys.stderr)
            for name, count, kinds in extra:
                shape = ", ".join(kinds) if kinds else "no content placeholders"
                print(f"  {name!r}: {count} content ({shape})", file=sys.stderr)
    # Partial or missing core layouts are a real finding, not a crash.
    core_keys = {s.key for s in CORE}
    return 1 if any(r.key in core_keys and not r.ok for r in results) else 0


if __name__ == "__main__":
    raise SystemExit(main())
