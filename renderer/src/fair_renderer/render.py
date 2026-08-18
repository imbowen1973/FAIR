"""The binding engine: session IR -> template-bound .pptx + JSON contracts.

Implements spec A.3 (binding mechanism) and A.4 (outputs).
"""

from __future__ import annotations

import copy
import json
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu

from .creation_id import derive_creation_id, inject_creation_id
from .layoutmap import LayoutBinding, LayoutMapError, load_layout_map, validate_against_template
from .mermaid import resolve_mermaid
from .parser import ListItem, Region, Session, Slide, parse_session
from .runs import write_runs
from .zipnorm import normalize_zip

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class RenderError(Exception):
    pass


def _clear_bullet_props(paragraph):
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    for tag in ("buChar", "buAutoNum", "buNone"):
        for el in pPr.findall(f"{{{A_NS}}}{tag}"):
            pPr.remove(el)
    return pPr


def _apply_auto_numbering(paragraph) -> None:
    """Mark a paragraph as auto-numbered (arabic, period) instead of bulleted."""
    pPr = _clear_bullet_props(paragraph)
    etree.SubElement(pPr, f"{{{A_NS}}}buAutoNum").set("type", "arabicPeriod")


def _suppress_bullet(paragraph) -> None:
    """Plain paragraph: no bullet glyph regardless of the layout's list style."""
    pPr = _clear_bullet_props(paragraph)
    etree.SubElement(pPr, f"{{{A_NS}}}buNone")


def _fill_list(placeholder, region: Region, numbered: bool) -> None:
    tf = placeholder.text_frame
    first = True

    def write_items(entries: list[ListItem], level: int) -> None:
        nonlocal first
        for item in entries:
            if first:
                paragraph = tf.paragraphs[0]
                first = False
            else:
                paragraph = tf.add_paragraph()
            paragraph.level = level
            write_runs(paragraph, item.text, color=item.color or region.color)
            if numbered:
                _apply_auto_numbering(paragraph)
            if item.children:
                write_items(item.children, level + 1)

    write_items(region.items, 0)


def _fill_plain(placeholder, region: Region, suppress_bullets: bool) -> None:
    tf = placeholder.text_frame
    lines = (region.text or "").splitlines() or [""]
    for i, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        write_runs(paragraph, line, color=region.color)
        if suppress_bullets:
            _suppress_bullet(paragraph)


def _fill_image(slide, placeholder, image_path: Path) -> None:
    """Insert an image into a region (spec A.3, image branch).

    Picture placeholders keep their <p:ph> binding via insert_picture().
    For non-picture placeholders, the image is placed at the placeholder's
    exact geometry and the empty placeholder removed — a documented
    exception to full placeholder binding (the shape carries no <p:ph>).
    """
    if not image_path.exists():
        raise RenderError(f"image not found: {image_path}")

    ph_type = placeholder.placeholder_format.type
    if ph_type == PP_PLACEHOLDER.PICTURE:
        return placeholder.insert_picture(str(image_path))

    # Contain-fit: scale to the largest size that fits the region without
    # distortion, centred in the placeholder's footprint.
    from PIL import Image

    with Image.open(image_path) as im:
        iw, ih = im.size
    box_w, box_h = placeholder.width, placeholder.height
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    left = Emu(placeholder.left + (box_w - w) // 2)
    top = Emu(placeholder.top + (box_h - h) // 2)
    pic = slide.shapes.add_picture(str(image_path), left, top, width=Emu(w), height=Emu(h))
    sp = placeholder._element
    sp.getparent().remove(sp)
    return pic


def _video_poster(build_dir: Path) -> Path:
    """Generate the default 16:9 poster (dark panel, play triangle)."""
    poster = build_dir / "video-poster.png"
    if poster.exists():
        return poster
    from PIL import Image, ImageDraw

    build_dir.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (1280, 720), (32, 38, 46))
    d = ImageDraw.Draw(img)
    cx, cy, r = 640, 360, 110
    d.ellipse([cx - r, cy - r, cx + r, cy + r], outline=(235, 238, 241), width=10)
    d.polygon([(cx - 35, cy - 62), (cx - 35, cy + 62), (cx + 72, cy)], fill=(235, 238, 241))
    img.save(poster)
    return poster


def _fill_video(slide, placeholder, region: Region, base_dir: Path, build_dir: Path) -> None:
    poster = base_dir / region.src if region.src else _video_poster(build_dir)
    pic = _fill_image(slide, placeholder, poster)
    pic.click_action.hyperlink.address = region.url


def _fill_region(slide, placeholder, region: Region, base_dir: Path, build_dir: Path) -> None:
    if region.type == "text":
        # Plain string region: emphasis-aware, keeps the layout's own
        # paragraph style (titles and heads are unbulleted by design).
        write_runs(placeholder.text_frame.paragraphs[0], region.text or "")
    elif region.type == "p":
        _fill_plain(placeholder, region, suppress_bullets=True)
    elif region.type in ("ul", "ol"):
        _fill_list(placeholder, region, numbered=(region.type == "ol"))
    elif region.type == "image":
        _fill_image(slide, placeholder, base_dir / region.src)
    elif region.type == "mermaid":
        image = resolve_mermaid(base_dir / region.src, build_dir)
        _fill_image(slide, placeholder, image)
    elif region.type == "video":
        _fill_video(slide, placeholder, region, base_dir, build_dir)
    else:  # pragma: no cover - parser guarantees the type set
        raise RenderError(f"unknown region type {region.type!r}")


def _render_slide(
    prs: Presentation,
    layouts_by_name: dict,
    binding: LayoutBinding,
    slide_src: Slide,
    session_id: str,
    base_dir: Path,
    build_dir: Path,
):
    layout = layouts_by_name[binding.layout_name]
    slide = prs.slides.add_slide(layout)

    placeholders_by_idx = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    for region in slide_src.regions:
        if region.name not in binding.regions:
            raise RenderError(
                f"slide {slide_src.id!r}: region {region.name!r} is not declared for "
                f"layout {binding.key!r}; expected one of {sorted(binding.regions)}"
            )
        idx = binding.regions[region.name]
        placeholder = placeholders_by_idx.get(idx)
        if placeholder is None:  # pragma: no cover - layout validation catches this
            raise RenderError(
                f"slide {slide_src.id!r}: placeholder idx {idx} missing on new slide"
            )
        _fill_region(slide, placeholder, region, base_dir, build_dir)

    if slide_src.notes:
        slide.notes_slide.notes_text_frame.text = slide_src.notes

    creation_id = derive_creation_id(session_id, slide_src.id)
    inject_creation_id(slide, creation_id)
    return slide, creation_id


def render_session(
    session_path: Path,
    template_path: Path,
    layout_map_path: Path,
    out_dir: Path,
) -> dict:
    """Render one session. Returns a summary dict with output paths."""
    session = parse_session(session_path)
    bindings = load_layout_map(layout_map_path)

    prs = Presentation(str(template_path))
    validate_against_template(bindings, prs, template_path)

    layouts_by_name = {
        layout.name: layout for master in prs.slide_masters for layout in master.slide_layouts
    }

    base_dir = session_path.parent
    build_dir = out_dir / "build"
    out_dir.mkdir(parents=True, exist_ok=True)

    session_id = session.session_id
    pptx_name = f"session-{session_id}.pptx"

    slide_id_map: dict[str, str] = {}
    index_entries: list[dict] = []

    for slide_src in session.slides:
        binding = bindings.get(slide_src.layout)
        if binding is None:
            raise RenderError(
                f"slide {slide_src.id!r}: layout {slide_src.layout!r} not in layout map; "
                f"available: {sorted(bindings)}"
            )
        slide, creation_id = _render_slide(
            prs, layouts_by_name, binding, slide_src, session_id, base_dir, build_dir
        )
        source_ref = f"{slide.slide_id}#{creation_id}"
        slide_id_map[slide_src.id] = source_ref

        title_region = next((r for r in slide_src.regions if r.name == "title"), None)
        plain_title = None
        if title_region and title_region.text:
            from .runs import parse_emphasis

            plain_title = "".join(r.text for r in parse_emphasis(title_region.text))
        index_entries.append(
            {
                "slideId": slide_src.id,
                "sessionId": session_id,
                "sourceRef": source_ref,
                "sourcePptx": pptx_name,
                "title": plain_title,
                "develops": slide_src.develops,
                "dok": slide_src.dok,
                # Carried so the pane can search narration, not just titles.
                "notes": slide_src.notes,
            }
        )

    pptx_path = out_dir / pptx_name
    prs.save(str(pptx_path))
    normalize_zip(pptx_path)

    map_path = out_dir / "slide-id-map.json"
    map_path.write_text(
        json.dumps(slide_id_map, indent=2, sort_keys=True) + "\n", encoding="utf-8"
    )

    index_path = out_dir / "index.json"
    index_doc = {
        "session": {
            key: value
            for key, value in session.frontmatter.items()
        },
        "slides": index_entries,
    }
    index_path.write_text(json.dumps(index_doc, indent=2) + "\n", encoding="utf-8")

    return {
        "pptx": pptx_path,
        "slide_id_map": map_path,
        "index": index_path,
        "slide_count": len(session.slides),
    }
