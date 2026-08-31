"""Load layout-map.yaml and validate it against the actual template.

The layout map is the explicit binding contract (spec A.2). Validation is
mandatory: if the template is re-edited in PowerPoint and placeholder
indices shift, the build must fail loudly rather than write content into
the wrong placeholder.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import yaml

from . import yamlio
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.presentation import Presentation


class LayoutMapError(Exception):
    """Raised when layout-map.yaml is malformed or disagrees with the template."""


@dataclass
class LayoutBinding:
    key: str  # key used in session.md, e.g. "Split"
    layout_name: str  # layout name inside template.pptx
    regions: dict[str, int]  # region name -> placeholder idx


def load_style(path: Path) -> dict:
    """Library-level style settings from the layout map's `_style` block.

    Only appearance the theme cannot express belongs here — today just
    `code_typeface`, because OOXML has major and minor font slots but no
    monospace one. It lives with the library so the renderer never names
    a font of its own.
    """
    try:
        data = yamlio.safe_load(path.read_text(encoding="utf-8"))
    except yaml.YAMLError as e:
        raise LayoutMapError(f"{path}: invalid YAML: {e}") from e
    style = (data or {}).get("_style") or {}
    if not isinstance(style, dict):
        raise LayoutMapError(f"{path}: '_style' must be a mapping")
    return style


def load_layout_map(path: Path) -> dict[str, LayoutBinding]:
    try:
        data = yamlio.safe_load(path.read_text(encoding="utf-8"))
    except yaml.YAMLError as e:
        raise LayoutMapError(f"{path}: invalid YAML: {e}") from e
    if not isinstance(data, dict) or not data:
        raise LayoutMapError(f"{path}: layout map must be a non-empty mapping")

    bindings: dict[str, LayoutBinding] = {}
    for key, entry in data.items():
        # Underscore keys are library settings, not layouts (see load_style).
        if key.startswith("_"):
            continue
        if not isinstance(entry, dict):
            raise LayoutMapError(f"{path}: entry {key!r} must be a mapping")
        layout_name = entry.get("layout_name")
        regions = entry.get("regions")
        if not isinstance(layout_name, str) or not layout_name:
            raise LayoutMapError(f"{path}: entry {key!r} missing 'layout_name'")
        if not isinstance(regions, dict):
            raise LayoutMapError(
                f"{path}: entry {key!r} missing 'regions' (use {{}} for chrome-only layouts)"
            )
        for rname, idx in regions.items():
            if not isinstance(idx, int) or idx < 0:
                raise LayoutMapError(
                    f"{path}: {key}.regions.{rname} must be a non-negative integer, got {idx!r}"
                )
        indices = list(regions.values())
        if len(indices) != len(set(indices)):
            raise LayoutMapError(f"{path}: entry {key!r} maps two regions to the same index")
        bindings[key] = LayoutBinding(key=key, layout_name=layout_name, regions=dict(regions))
    return bindings


# The name a slide uses for "the content placeholder", whichever one
# this layout calls it.
MAIN_CONTENT_REGION = "full"

# What PowerPoint will let you drop a picture into. A BODY placeholder
# takes text and nothing else, which is why a section header's subtitle
# is not a candidate however alone it is on the slide.
CONTENT_TYPES = {PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.PICTURE}


def _main_content(binding: LayoutBinding, layout) -> str | None:
    """The one region of this layout that holds the slide's content.

    A layout map names the same PowerPoint placeholder differently in
    every layout — `full` in Title and Content, `picture` in Picture with
    Caption, `content` in Content with Caption — so a slide could not be
    moved between them, and content written for one was undeclared in
    the next. That is a naming accident, not a difference in the
    template: all three are the placeholder a picture goes into.

    So where a layout has exactly one of those, it answers to `full` as
    well as to its own name. Where it has two — Comparison, Split, Cards
    — there is a real choice to make and this returns None, because
    guessing which column somebody meant is worse than asking.
    """
    types = {ph.placeholder_format.idx: ph.placeholder_format.type for ph in layout.placeholders}
    holders = [r for r, idx in binding.regions.items() if types.get(idx) in CONTENT_TYPES]
    return holders[0] if len(holders) == 1 else None


def apply_content_aliases(bindings: dict[str, LayoutBinding], prs: Presentation) -> None:
    """Let every single-content layout answer to `full`.

    Applied to the bindings rather than written into layout-map.yaml, so
    libraries that already exist gain it without regenerating anything.
    A map that declares `full` itself is left alone.
    """
    layouts_by_name = {
        layout.name: layout for master in prs.slide_masters for layout in master.slide_layouts
    }
    for binding in bindings.values():
        if MAIN_CONTENT_REGION in binding.regions:
            continue
        layout = layouts_by_name.get(binding.layout_name)
        if layout is None:
            continue
        region = _main_content(binding, layout)
        if region is not None:
            binding.regions[MAIN_CONTENT_REGION] = binding.regions[region]


def validate_against_template(
    bindings: dict[str, LayoutBinding], prs: Presentation, template_path: Path
) -> None:
    """Fail fast if any mapped layout or placeholder index is absent."""
    layouts_by_name = {
        layout.name: layout for master in prs.slide_masters for layout in master.slide_layouts
    }
    errors: list[str] = []
    for binding in bindings.values():
        layout = layouts_by_name.get(binding.layout_name)
        if layout is None:
            errors.append(
                f"layout {binding.layout_name!r} (key {binding.key!r}) not found in template; "
                f"template has: {sorted(layouts_by_name)}"
            )
            continue
        available = {ph.placeholder_format.idx for ph in layout.placeholders}
        for rname, idx in binding.regions.items():
            if idx not in available:
                errors.append(
                    f"{binding.key}.{rname} -> idx {idx}, but layout "
                    f"{binding.layout_name!r} only has placeholder indices {sorted(available)}"
                )
    if errors:
        raise LayoutMapError(
            f"layout map does not match template {template_path}:\n  " + "\n  ".join(errors)
        )
