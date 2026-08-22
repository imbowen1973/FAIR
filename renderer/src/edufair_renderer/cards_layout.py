"""Inject the Cards layout into an existing PowerPoint template.

Cards is an Extended-profile layout: no stock template has it, and asking
a designer to hand-build nine placeholders at exact indices would defeat
"any standard template works". So we add it to whatever template they
brought.

Nothing here carries a colour value. The tabs fill from `accent1`-`accent4`,
the panels from `bg2`, the tab text from `lt1` — theme slots that OOXML
guarantees every theme defines. The cards therefore adopt the target
template's palette, and follow it again if that template is re-themed.

python-pptx has no API for adding a slide layout, so this works at the
OPC level: clone a spare layout part, strip it, inject the placeholders,
and register the new part on the master.
"""

from __future__ import annotations

import copy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.slide import SlideLayoutPart

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NSMAP = f'xmlns:a="{A_NS}" xmlns:p="{P_NS}"'
EMU_PER_IN = 914400

# Layouts we are willing to clone: title plus chrome, no content to strip.
DONORS = ("Title Only", "TitleOnly", "Blank")

# p:sldLayoutId/@id must be >= 2147483648 per ST_SlideLayoutId.
_SLD_LAYOUT_ID_BASE = 2147483649


class CardsError(Exception):
    pass


CARD_TAB_XML = """\
<p:sp {ns}>
  <p:nvSpPr>
    <p:cNvPr id="{id}" name="Card {n} Tab"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="body" sz="quarter" idx="{idx}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="round2SameRect">
      <a:avLst><a:gd name="adj1" fmla="val 25000"/><a:gd name="adj2" fmla="val 0"/></a:avLst>
    </a:prstGeom>
    <a:solidFill><a:schemeClr val="{accent}"/></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr" lIns="91440" rIns="91440" tIns="0" bIns="0"/>
    <a:lstStyle>
      <a:lvl1pPr algn="ctr" marL="0" indent="0"><a:buNone/>
        <a:defRPr sz="{tab_pt}" b="1"><a:solidFill><a:schemeClr val="lt1"/></a:solidFill></a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr/></a:p>
  </p:txBody>
</p:sp>
"""

CARD_BODY_XML = """\
<p:sp {ns}>
  <p:nvSpPr>
    <p:cNvPr id="{id}" name="Card {n} Body"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="body" sz="quarter" idx="{idx}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:schemeClr val="bg2"/></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr lIns="91440" rIns="91440" tIns="82296" bIns="82296">
      <a:normAutofit/>
    </a:bodyPr>
    <a:lstStyle>
      <a:lvl1pPr><a:defRPr sz="{body_pt}"/></a:lvl1pPr>
      <a:lvl2pPr><a:defRPr sz="{sub_pt}"/></a:lvl2pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr/></a:p>
  </p:txBody>
</p:sp>
"""


def card_geometry(slide_width: int, slide_height: int) -> dict:
    """Card positions for a slide of any size.

    Derived from both dimensions: fixed inch constants would overflow a
    short widescreen slide (10 x 5.63) while looking sparse on a tall one.
    """
    # Integer fractions, not float inches: these reproduce the original
    # 10 x 7.5in geometry exactly (so regenerating the committed template
    # is byte-stable) while still scaling to any slide size.
    margin = slide_width * 33 // 1000
    gap = slide_width * 17 // 1000
    card_w = (slide_width - 2 * margin - 3 * gap) // 4

    top = slide_height * 180 // 750       # below the title band
    tab_h = slide_height * 55 // 750
    body_h = slide_height * 475 // 750

    # Type scales with the slide so cards stay readable at any size.
    scale = slide_height / (7.5 * EMU_PER_IN)
    return {
        "margin": margin,
        "gap": gap,
        "card_w": card_w,
        "top": top,
        "tab_h": tab_h,
        "body_h": body_h,
        "tab_pt": max(800, int(1500 * scale)),
        "body_pt": max(700, int(1300 * scale)),
        "sub_pt": max(600, int(1100 * scale)),
    }


def build_card_shapes(slide_width: int, slide_height: int) -> list:
    """The eight placeholders: tabs at idx 1-4, panels at idx 5-8."""
    g = card_geometry(slide_width, slide_height)
    shapes = []
    for n in range(1, 5):
        x = g["margin"] + (n - 1) * (g["card_w"] + g["gap"])
        shapes.append(
            etree.fromstring(
                CARD_TAB_XML.format(
                    ns=NSMAP, id=20 + n, n=n, idx=n,
                    x=x, y=g["top"], cx=g["card_w"], cy=g["tab_h"],
                    accent=f"accent{n}", tab_pt=g["tab_pt"],
                )
            )
        )
        shapes.append(
            etree.fromstring(
                CARD_BODY_XML.format(
                    ns=NSMAP, id=24 + n, n=n, idx=4 + n,
                    x=x, y=g["top"] + g["tab_h"], cx=g["card_w"], cy=g["body_h"],
                    body_pt=g["body_pt"], sub_pt=g["sub_pt"],
                )
            )
        )
    return shapes


def _strip_content_placeholders(layout_el) -> None:
    """Keep the title (idx 0) and chrome (idx >= 10); drop the rest."""
    sp_tree = layout_el.find(f"{{{P_NS}}}cSld/{{{P_NS}}}spTree")
    for sp in list(sp_tree.findall(f"{{{P_NS}}}sp")):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None:
            continue
        idx = int(ph.get("idx") or 0)
        if 0 < idx < 10:
            sp_tree.remove(sp)


def _pick_donor(master):
    for name in DONORS:
        for layout in master.slide_layouts:
            if layout.name.strip().lower() == name.strip().lower():
                return layout
    # Any layout will do once its content placeholders are stripped.
    if len(master.slide_layouts) == 0:
        raise CardsError("template has no slide layouts to clone")
    return master.slide_layouts[0]


def has_cards(prs: Presentation) -> str | None:
    """The name of an existing Cards-shaped layout, or None.

    Checked by name and by shape, so a template that already has a
    quad-panel layout is reused rather than growing a second one.
    """
    from .template_inspect import _find_cards

    layouts = {
        layout.name: layout
        for master in prs.slide_masters
        for layout in master.slide_layouts
    }
    for name in layouts:
        if name.strip().lower() == "cards":
            return name
    name, _, _ = _find_cards(layouts, set())
    return name


def inject_cards(template: Path, out: Path, force: bool = False) -> dict:
    """Add a Cards layout to `template`, writing the result to `out`.

    Returns {"added": bool, "layout": name, "donor": name|None}.
    """
    prs = Presentation(str(template))

    existing = None if force else has_cards(prs)
    if existing:
        prs.save(str(out))
        return {"added": False, "layout": existing, "donor": None}

    master = prs.slide_masters[0]
    donor = _pick_donor(master)

    # Clone rather than repurpose: the donor is someone else's layout and
    # must survive untouched.
    layout_el = copy.deepcopy(donor.element)
    layout_el.find(f"{{{P_NS}}}cSld").set("name", "Cards")
    _strip_content_placeholders(layout_el)

    sp_tree = layout_el.find(f"{{{P_NS}}}cSld/{{{P_NS}}}spTree")
    for shape in build_card_shapes(prs.slide_width, prs.slide_height):
        sp_tree.append(shape)

    package = prs.part.package
    partname = package.next_partname("/ppt/slideLayouts/slideLayout%d.xml")
    blob = etree.tostring(layout_el, xml_declaration=True, encoding="UTF-8", standalone=True)
    part = SlideLayoutPart.load(partname, donor.part.content_type, package, blob)

    # A layout must point at its master, and the master must list it.
    part.relate_to(master.part, RT.SLIDE_MASTER)
    rId = master.part.relate_to(part, RT.SLIDE_LAYOUT)

    id_lst = master.element.get_or_add_sldLayoutIdLst()
    entry = id_lst._add_sldLayoutId()
    entry.set("id", str(_SLD_LAYOUT_ID_BASE + len(id_lst.sldLayoutId_lst)))
    entry.set(f"{{{R_NS}}}id", rId)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return {"added": True, "layout": "Cards", "donor": donor.name}
