"""Attribution: visible, locked, and machine-readable — per slide.

Three layers, injected at render time (rendering is the only path from
content to deck, so attribution cannot be forgotten):

1. A visible stamp — small text (and optional logo) in a corner of
   every slide, as shapes carrying OOXML locks (noSelect etc.): they
   cannot be clicked, moved, or deleted in PowerPoint's UI. Removal
   requires editing the file's XML — deliberate, not casual.
2. A line appended to the speaker notes.
3. A provenance extension in each slide's XML (beside the creationId):
   origin, licence, session and slide ids. PowerPoint preserves unknown
   extensions through edits, re-theming and cross-deck copies, so this
   survives even after a determined user strips the visible marks —
   which is what makes stripping provable. `edufair-audit` reads it back.

Config: attribution.yaml at the library root (next to sessions/):

    text: "CC-BY 4.0 · FAIR Consortium"
    logo: branding/logo.png      # optional, relative to this file
    corner: bottom-right         # or bottom-left
    opacity: 0.45                # optional, 0-1; the stamp text
    scale: 1.0                   # optional, multiplies the text size
    logo_height_cm: 1.0          # optional, overrides the logo's height
    logo_opacity: 1.0            # optional, defaults to opaque

The logo is sized and faded independently of the text on purpose. An
EU emblem carries rules the wording does not: a 1 cm minimum height,
and no alteration — so it must be able to sit at full strength beside
text that is small and faded.

No file -> no visible stamp; the provenance extension is written
regardless, from the session's frontmatter (dc.creator, dc.license).
"""

from __future__ import annotations

from pathlib import Path

import yaml
from lxml import etree
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
FAIR_NS = "urn:fair:attribution:1"
# Registered extension URI for the FAIR provenance block (fixed GUID).
ATTRIBUTION_EXT_URI = "{9D2C8E4A-5B31-4F60-A0C7-2D8E13AB0042}"

CORNERS = {"bottom-right", "bottom-left"}
_EMU_IN = 914400


class AttributionError(Exception):
    pass


def load_attribution(path: Path) -> dict:
    data = yaml.safe_load(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict) or not isinstance(data.get("text"), str):
        raise AttributionError(f"{path}: attribution.yaml needs a 'text' string")
    corner = data.get("corner", "bottom-right")
    if corner not in CORNERS:
        raise AttributionError(f"{path}: corner must be one of {sorted(CORNERS)}")
    logo = data.get("logo")
    if logo is not None:
        logo_path = (path.parent / logo).resolve()
        if not logo_path.exists():
            raise AttributionError(f"{path}: logo not found: {logo}")
        data["_logo_path"] = logo_path
    data["corner"] = corner

    opacity = data.get("opacity", 1.0)
    if not isinstance(opacity, (int, float)) or not 0 < opacity <= 1:
        raise AttributionError(f"{path}: opacity must be a number in (0, 1]")
    data["opacity"] = float(opacity)

    scale = data.get("scale", 1.0)
    if not isinstance(scale, (int, float)) or not 0 < scale <= 4:
        raise AttributionError(f"{path}: scale must be a number in (0, 4]")
    data["scale"] = float(scale)

    # The logo is sized and faded on its own terms: an EU emblem has a
    # 1 cm minimum height and must not be altered, so it cannot be tied
    # to whatever keeps the wording unobtrusive.
    height_cm = data.get("logo_height_cm")
    if height_cm is not None:
        if not isinstance(height_cm, (int, float)) or not 0 < height_cm <= 10:
            raise AttributionError(f"{path}: logo_height_cm must be a number in (0, 10]")
        data["logo_height_cm"] = float(height_cm)

    logo_opacity = data.get("logo_opacity", 1.0)
    if not isinstance(logo_opacity, (int, float)) or not 0 < logo_opacity <= 1:
        raise AttributionError(f"{path}: logo_opacity must be a number in (0, 1]")
    data["logo_opacity"] = float(logo_opacity)
    return data


def _apply_alpha(color_element, opacity: float) -> None:
    """DrawingML alpha: 0-100000, on a colour or a blip fix."""
    if opacity >= 1:
        return
    alpha = etree.SubElement(color_element, f"{{{A_NS}}}alpha")
    alpha.set("val", str(int(round(opacity * 100000))))


def _lock(nv_pr_element, tag: str, attrs: tuple[str, ...]) -> None:
    # python-pptx may already have created the locks element (e.g.
    # picLocks noChangeAspect); merge into it rather than duplicating.
    locks = nv_pr_element.find(f"{{{A_NS}}}{tag}")
    if locks is None:
        locks = etree.SubElement(nv_pr_element, f"{{{A_NS}}}{tag}")
    for attr in attrs:
        locks.set(attr, "1")


def _stamped_logo(logo_path: Path, build_dir: Path, meta: dict) -> Path:
    """Copy the logo with attribution baked into the PNG's own metadata."""
    from PIL import Image
    from PIL.PngImagePlugin import PngInfo

    out = build_dir / "attribution-logo.png"
    if out.exists():
        return out
    build_dir.mkdir(parents=True, exist_ok=True)
    info = PngInfo()
    for key in ("creator", "license", "text"):
        if meta.get(key):
            info.add_text(f"fair:{key}", str(meta[key]))
    with Image.open(logo_path) as im:
        im.load()
        im.save(out, format="PNG", pnginfo=info)
    return out


def stamp_slide(slide, prs, config: dict, build_dir: Path, meta: dict) -> None:
    """Layer 1: the visible, locked stamp."""
    text = config["text"]
    corner = config["corner"]
    opacity = config.get("opacity", 1.0)
    scale = config.get("scale", 1.0)
    logo_opacity = config.get("logo_opacity", 1.0)
    height_cm = config.get("logo_height_cm")
    margin = Emu(int(0.12 * _EMU_IN))
    box_h = Emu(int(0.24 * scale * _EMU_IN))
    box_w = Emu(int(3.4 * scale * _EMU_IN))
    logo_h = Emu(
        int((height_cm / 2.54 if height_cm else 0.26 * scale) * _EMU_IN)
    )
    # Sit the wording on the logo's baseline band when the logo is taller.
    top = prs.slide_height - max(box_h, logo_h) - margin + (max(box_h, logo_h) - box_h)

    logo_path = config.get("_logo_path")
    logo_w = Emu(0)
    if logo_path:
        from PIL import Image

        stamped = _stamped_logo(logo_path, build_dir, {**meta, "text": text})
        with Image.open(stamped) as im:
            ratio = im.width / im.height
        logo_w = Emu(int(int(logo_h) * ratio))
        logo_left = (
            prs.slide_width - margin - logo_w
            if corner == "bottom-right"
            else margin
        )
        pic = slide.shapes.add_picture(
            str(stamped), logo_left, prs.slide_height - logo_h - margin,
            width=logo_w, height=logo_h,
        )
        c_nv_pic_pr = pic._element.find(f"{{{P_NS}}}nvPicPr/{{{P_NS}}}cNvPicPr")
        _lock(
            c_nv_pic_pr,
            "picLocks",
            ("noSelect", "noMove", "noResize", "noGrp", "noChangeAspect"),
        )
        # alphaModFix inside the blip is how DrawingML fades a picture;
        # it survives re-theming because it is part of the fill, not a style.
        if logo_opacity < 1:
            blip = pic._element.blipFill.find(f"{{{A_NS}}}blip")
            fix = etree.SubElement(blip, f"{{{A_NS}}}alphaModFix")
            fix.set("amt", str(int(round(logo_opacity * 100000))))

    gap = Emu(int(0.06 * _EMU_IN)) if logo_path else Emu(0)
    if corner == "bottom-right":
        left = prs.slide_width - margin - logo_w - gap - box_w
        align = PP_ALIGN.RIGHT
    else:
        left = margin + logo_w + gap
        align = PP_ALIGN.LEFT

    tb = slide.shapes.add_textbox(left, top, box_w, box_h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(max(1, round(8 * scale)))
    from pptx.dml.color import RGBColor

    run.font.color.rgb = RGBColor(0x8A, 0x8A, 0x8A)
    # Same fade on the text, so mark and wording read as one stamp.
    _apply_alpha(run.font.color._xFill.find(f"{{{A_NS}}}srgbClr"), opacity)
    _lock(
        tb._element.nvSpPr.cNvSpPr,
        "spLocks",
        ("noSelect", "noMove", "noResize", "noGrp", "noTextEdit"),
    )


def append_notes_line(slide, text: str, session_id: str, slide_id: str) -> None:
    """Layer 2: attribution in the speaker notes."""
    tf = slide.notes_slide.notes_text_frame
    paragraph = tf.paragraphs[0] if not tf.paragraphs[0].runs and len(tf.paragraphs) == 1 else tf.add_paragraph()
    run = paragraph.add_run()
    run.text = f"Source: {text} · {session_id}/{slide_id}"


def inject_provenance(slide, meta: dict) -> None:
    """Layer 3: machine-readable provenance beside the creationId."""
    sld = slide._element
    ext_lst = sld.find(f"{{{P_NS}}}extLst")
    if ext_lst is None:
        ext_lst = etree.SubElement(sld, f"{{{P_NS}}}extLst")
    ext = etree.SubElement(ext_lst, f"{{{P_NS}}}ext")
    ext.set("uri", ATTRIBUTION_EXT_URI)
    node = etree.SubElement(ext, f"{{{FAIR_NS}}}attribution", nsmap={"fair": FAIR_NS})
    for key in ("creator", "license", "session", "slide", "version"):
        if meta.get(key):
            node.set(key, str(meta[key]))
