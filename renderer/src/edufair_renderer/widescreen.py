"""Widen a 4:3 template to 16:9, keeping the design's proportions.

A library made before the seed template was widened starts 4:3, and there
is no way to tell from inside the workbench: the canvas draws the aspect
the template declares, so a 4:3 outline is an accurate report of a 4:3
deck rather than a bug in the outline. Changing the canvas would make it
lie while the rendered .pptx stayed 4:3.

So the template is what changes. The slide becomes 12192000 EMU wide --
the same 6858000 tall -- and every horizontal offset and extent is scaled
by the same 4/3, so a design keeps its proportions instead of sitting in
a column beside an empty band.

Done on the XML, once per element. python-pptx exposes the same offset
through more than one wrapper, and scaling through the wrappers visited
some shapes twice: 1.333 x 1.333 is 1.778, which is how placeholders came
to hang off the right-hand edge the first time this was attempted.

Vertical values are untouched. The height does not change, so scaling
them would be wrong, and a shape that filled the slide top to bottom
still does.
"""

from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

from pptx import Presentation

WIDESCREEN_EMU = 12192000
STANDARD_EMU = 9144000

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Every attribute that measures across the slide. `chOff`/`chExt` are the
# child coordinate space of a group and live in the same units, so a group
# whose children were not scaled with it would tear itself apart.
_HORIZONTAL = (
    (f"{_A}off", "x"),
    (f"{_A}ext", "cx"),
    (f"{_A}chOff", "x"),
    (f"{_A}chExt", "cx"),
)


def _scale_part(root, factor: float) -> int:
    """Scale every horizontal measurement in one XML part. Returns how many."""
    touched = 0
    for tag, attr in _HORIZONTAL:
        for node in root.iter(tag):
            value = node.get(attr)
            if value is None:
                continue
            try:
                node.set(attr, str(round(int(value) * factor)))
            except ValueError:
                continue  # a value we do not understand is one to leave alone
            touched += 1
    return touched


def widen(path: Path, out: Path | None = None, width: int = WIDESCREEN_EMU) -> dict:
    """Widen `path` to `width`, writing to `out` (or over `path`).

    Returns what changed, so a caller can say so rather than claiming it.
    """
    prs = Presentation(str(path))
    was = prs.slide_width
    if was == width:
        return {"changed": False, "from": was, "to": width, "elements": 0, "parts": 0}
    factor = width / was

    parts = 0
    touched = 0
    seen = set()
    for master in prs.slide_masters:
        for part in (master, *master.slide_layouts):
            # A layout belongs to one master, but a master can be reached
            # more than once. Scaling a part twice squares the factor,
            # which is the fault this exists to avoid.
            if id(part.element) in seen:
                continue
            seen.add(id(part.element))
            touched += _scale_part(part.element, factor)
            parts += 1
    for slide in prs.slides:
        if id(slide.element) in seen:
            continue
        seen.add(id(slide.element))
        touched += _scale_part(slide.element, factor)
        parts += 1

    prs.slide_width = width
    prs.save(str(out or path))
    return {"changed": True, "from": was, "to": width, "elements": touched, "parts": parts}


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(
        prog="edufair-widescreen",
        description=(
            "Widen a 4:3 template to 16:9, scaling every horizontal offset and "
            "extent so the design keeps its proportions. Re-run edufair-template "
            "afterwards: layout-geometry.json records where the placeholders were."
        ),
    )
    ap.add_argument("template", type=Path, help="the template.pptx to widen")
    ap.add_argument("--out", type=Path, default=None, help="write here instead of in place")
    ap.add_argument(
        "--backup",
        action="store_true",
        help="keep the original beside it as template.pptx.4x3",
    )
    args = ap.parse_args(argv)

    if not args.template.is_file():
        print(f"no such template: {args.template}", file=sys.stderr)
        return 2
    if args.backup and args.out is None:
        shutil.copy2(args.template, args.template.with_suffix(".pptx.4x3"))

    result = widen(args.template, args.out)
    if not result["changed"]:
        print(f"{args.template} is already {result['to']} EMU wide — nothing to do.")
        return 0
    print(
        f"{args.template}: {result['from']} -> {result['to']} EMU "
        f"({result['elements']} measurements across {result['parts']} parts). "
        "Now re-run edufair-template to rewrite layout-geometry.json."
    )
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
