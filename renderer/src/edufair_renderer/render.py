"""The binding engine: session IR -> template-bound .pptx + JSON contracts.

Implements spec A.3 (binding mechanism) and A.4 (outputs).
"""

from __future__ import annotations

import copy
import json
from pathlib import Path

import yaml

from . import yamlio
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu

from .attribution import (
    append_notes_line,
    inject_provenance,
    load_attribution,
    stamp_slide,
)
from .creation_id import derive_creation_id, inject_creation_id
from .layoutmap import (
    LayoutBinding,
    LayoutMapError,
    load_layout_map,
    load_style,
    validate_against_template,
)
from .mermaid import resolve_mermaid
from .outcomes import OUTCOMES_FILE, derive_develops, fill_roles, load_outcomes
from .parser import ListItem, Region, Session, Slide, parse_session
from .runs import DEFAULT_CODE_TYPEFACE, plain_text, write_runs
from .zipnorm import normalize_zip

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class RenderError(Exception):
    pass


def _clear_bullet_props(paragraph):
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    for tag in ("buChar", "buAutoNum", "buNone"):
        for el in pPr.findall(f"{{{A_NS}}}{tag}"):
            pPr.remove(el)
    return pPr


def _apply_auto_numbering(paragraph) -> None:
    """Mark a paragraph as auto-numbered (arabic, period) instead of bulleted."""
    pPr = _clear_bullet_props(paragraph)
    etree.SubElement(pPr, f"{{{A_NS}}}buAutoNum").set("type", "arabicPeriod")


def _suppress_bullet(paragraph) -> None:
    """Plain paragraph: no bullet glyph regardless of the layout's list style."""
    pPr = _clear_bullet_props(paragraph)
    etree.SubElement(pPr, f"{{{A_NS}}}buNone")


def _bullet_sources(placeholder):
    """The list-style elements a body placeholder inherits from, in order.

    The placeholder's own lstStyle, then the layout's matching
    placeholder, then the master's bodyStyle. Nothing else: a first
    attempt searched the whole master, which has bullets defined
    *somewhere* in every real template, so it always concluded the
    template had an opinion and never supplied one.
    """
    out = []
    body = placeholder._element.find(f"{{{A_NS}}}txBody")
    if body is not None:
        out.append(body.find(f"{{{A_NS}}}lstStyle"))
    try:
        idx = placeholder.placeholder_format.idx
        layout = placeholder.part.slide_layout
        for shape in layout.placeholders:
            if shape.placeholder_format.idx == idx:
                shape_body = shape._element.find(f"{{{A_NS}}}txBody")
                if shape_body is not None:
                    out.append(shape_body.find(f"{{{A_NS}}}lstStyle"))
                break
        master = layout.slide_master._element
        txStyles = master.find(f"{{{A_NS}}}txStyles")
        if txStyles is not None:
            out.append(txStyles.find(f"{{{A_NS}}}bodyStyle"))
    except Exception:  # not a placeholder, or no layout behind it
        pass
    return [o for o in out if o is not None]


def _defines_a_bullet(placeholder, level: int) -> bool:
    """Whether the template puts a bullet on paragraphs at this level.

    A `ul` that renders without bullets is simply wrong: the content says
    the list is bulleted and the output does not show it. The glyph and
    its size stay the template's, so this only asks whether the template
    has an opinion at all -- and one is supplied when it has none.

    The example template's master has no `bodyStyle` element whatsoever,
    so every bulleted list came out as plain paragraphs.
    """
    want = f"{{{A_NS}}}lvl{level + 1}pPr"
    for style in _bullet_sources(placeholder):
        level_props = style.find(want)
        if level_props is None:
            continue
        for tag in ("buChar", "buAutoNum"):
            if level_props.find(f"{{{A_NS}}}{tag}") is not None:
                return True
        # An explicit buNone at this level is the template saying no.
        if level_props.find(f"{{{A_NS}}}buNone") is not None:
            return True
    return False


ALIGNMENTS = {
    "left": "l",
    "center": "ctr",
    "right": "r",
    "justify": "just",
}


def _apply_alignment(paragraph, align: str) -> None:
    """Align one paragraph, leaving the rest of the placeholder alone.

    A centred heading above left-aligned points is one placeholder in
    every real deck, and alignment was a property of the whole region.
    """
    value = ALIGNMENTS.get(align)
    if not value:
        return
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("algn", value)


def _apply_bullet(paragraph, char: str = "•") -> None:
    """Put a bullet on a paragraph the template left bare."""
    pPr = _clear_bullet_props(paragraph)
    element = etree.SubElement(pPr, f"{{{A_NS}}}buChar")
    element.set("char", char)


def _fill_list(placeholder, region: Region, numbered: bool, code_typeface: str) -> None:
    tf = placeholder.text_frame
    first = True
    # Per level, asked once: a template can define bullets for the first
    # level and leave the nested ones bare, and walking the inheritance
    # chain for every paragraph would give the same answer each time.
    bare_level = {}

    def needs_bullet(level: int) -> bool:
        if level not in bare_level:
            bare_level[level] = not _defines_a_bullet(placeholder, level)
        return bare_level[level]

    def write_items(entries: list[ListItem], level: int) -> None:
        nonlocal first
        for item in entries:
            if first:
                paragraph = tf.paragraphs[0]
                first = False
            else:
                paragraph = tf.add_paragraph()
            paragraph.level = level
            write_runs(
                paragraph,
                item.text,
                color=item.color or region.color,
                code_typeface=code_typeface,
            )
            # Per line. The region says what the list is; a line may say
            # otherwise, which is how one placeholder holds a lead-in,
            # then the points, then the sentence that closes them.
            marker = item.marker or ("number" if numbered else "bullet")
            if marker == "none":
                _suppress_bullet(paragraph)
            elif marker == "number":
                _apply_auto_numbering(paragraph)
            elif needs_bullet(level):
                _apply_bullet(paragraph)
            if item.align:
                _apply_alignment(paragraph, item.align)
            if item.children:
                write_items(item.children, level + 1)

    write_items(region.items, 0)


def _fill_plain(placeholder, region: Region, suppress_bullets: bool, code_typeface: str) -> None:
    tf = placeholder.text_frame
    lines = (region.text or "").splitlines() or [""]
    for i, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        write_runs(paragraph, line, color=region.color, code_typeface=code_typeface)
        if suppress_bullets:
            _suppress_bullet(paragraph)


# How a picture meets its frame. A portrait photograph in a landscape
# placeholder needs a decision, and it is not one the template can make:
# the template knows the shape of the hole, only the author knows whether
# this particular picture may lose its edges.
IMAGE_FITS = ("cover", "contain", "width", "height")
DEFAULT_IMAGE_FIT = "cover"


def _fit_geometry(frame_w, frame_h, img_w, img_h, fit):
    """Where a picture sits in its frame, and what is cropped off it.

    Returns `(width, height, crop_x, crop_y)` — the shape's size, and the
    fraction to take off *each* side of the source. One of the crops is
    always zero: a picture is never cropped in both directions at once.
    """
    if not img_w or not img_h:
        return frame_w, frame_h, 0.0, 0.0

    by_width = frame_w / img_w
    by_height = frame_h / img_h
    scale = {
        "cover": max(by_width, by_height),
        "contain": min(by_width, by_height),
        "width": by_width,
        "height": by_height,
    }.get(fit, max(by_width, by_height))

    drawn_w, drawn_h = img_w * scale, img_h * scale

    # Wider than the frame: keep the frame and crop the sides. Narrower:
    # shrink the shape and let it sit centred in the space.
    if drawn_w > frame_w + 1:
        crop_x = (1 - frame_w / drawn_w) / 2
        width = frame_w
    else:
        crop_x = 0.0
        width = int(round(drawn_w))

    if drawn_h > frame_h + 1:
        crop_y = (1 - frame_h / drawn_h) / 2
        height = frame_h
    else:
        crop_y = 0.0
        height = int(round(drawn_h))

    return width, height, crop_x, crop_y


def _fill_image(slide, placeholder, image_path: Path, fit: str = DEFAULT_IMAGE_FIT) -> None:
    """Insert an image into a region (spec A.3, image branch).

    Picture placeholders keep their <p:ph> binding via insert_picture().
    For non-picture placeholders, the image is placed at the placeholder's
    exact geometry and the empty placeholder removed — a documented
    exception to full placeholder binding (the shape carries no <p:ph>).

    Both kinds obey `fit` the same way. They used not to: a picture
    placeholder cropped to fill while any other placeholder fitted
    inside, so the same picture behaved differently depending on which
    layout it landed in, and neither could be told otherwise.
    """
    if not image_path.exists():
        raise RenderError(f"image not found: {image_path}")

    from PIL import Image

    with Image.open(image_path) as im:
        img_w, img_h = im.size

    frame_w, frame_h = placeholder.width, placeholder.height
    left, top = placeholder.left, placeholder.top
    width, height, crop_x, crop_y = _fit_geometry(
        frame_w, frame_h, img_w, img_h, fit
    )

    ph_type = placeholder.placeholder_format.type
    if ph_type == PP_PLACEHOLDER.PICTURE:
        # insert_picture keeps the <p:ph> and fills the frame, cropping;
        # the crops and the size are then set to whatever `fit` asked for.
        pic = placeholder.insert_picture(str(image_path))
    else:
        pic = slide.shapes.add_picture(str(image_path), Emu(left), Emu(top))
        sp = placeholder._element
        sp.getparent().remove(sp)

    pic.crop_left = crop_x
    pic.crop_right = crop_x
    pic.crop_top = crop_y
    pic.crop_bottom = crop_y
    pic.width = Emu(int(width))
    pic.height = Emu(int(height))
    # Centred in the frame, which matters whenever the picture does not
    # fill it: an image pinned to the top-left of its hole looks wrong.
    pic.left = Emu(int(left + (frame_w - width) / 2))
    pic.top = Emu(int(top + (frame_h - height) / 2))
    return pic


def _video_poster(build_dir: Path) -> Path:
    """Generate the default 16:9 poster (dark panel, play triangle)."""
    poster = build_dir / "video-poster.png"
    if poster.exists():
        return poster
    from PIL import Image, ImageDraw

    build_dir.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (1280, 720), (32, 38, 46))
    d = ImageDraw.Draw(img)
    cx, cy, r = 640, 360, 110
    d.ellipse([cx - r, cy - r, cx + r, cy + r], outline=(235, 238, 241), width=10)
    d.polygon([(cx - 35, cy - 62), (cx - 35, cy + 62), (cx + 72, cy)], fill=(235, 238, 241))
    img.save(poster)
    return poster


def _library_root(base_dir: Path) -> Path | None:
    """The library a block belongs to, if this is a library at all.

    A block lives at <root>/blocks/<id>, so the root is two up. A flat
    sessions directory has no root and gets None.
    """
    parent = base_dir.parent
    return parent.parent if parent.name == "blocks" else None


def _asset_path(base_dir: Path, src: str) -> Path:
    """Where the media a slide names actually is.

    A slide names its media relative to its own block, which is what an
    upload writes: `media/consortium.jpg`. But the workbench's picker
    offers every image in the library and writes the repo-relative path
    for one it did not just upload, so the block got counted twice:

        blocks/01-intro/blocks/01-intro/media/consortium.jpg

    Those slides are committed and the canvas draws them perfectly --
    it accepts both spellings -- so refusing them here fails a deck for
    a difference the author cannot see. A path that resolves from the
    library root is accepted as well.

    Which also makes reuse across blocks expressible, and it has to be:
    the picker lists the whole library, so it was already being asked
    for.

    The block-relative path is returned when neither exists, so the
    error names the place a reader would look first.
    """
    direct = base_dir / src
    if direct.exists():
        return direct
    root = _library_root(base_dir)
    if root is not None:
        from_root = root / src
        if from_root.exists():
            return from_root
    return direct


def _fill_video(slide, placeholder, region: Region, base_dir: Path, build_dir: Path) -> None:
    poster = _asset_path(base_dir, region.src) if region.src else _video_poster(build_dir)
    pic = _fill_image(slide, placeholder, poster, region.fit or DEFAULT_IMAGE_FIT)
    pic.click_action.hyperlink.address = region.url


def _shrink_on_overflow(placeholder) -> None:
    """Let PowerPoint scale text that outgrows its placeholder.

    Without this a long list simply spills past the bottom of its box and
    collides with whatever is beneath it — the placeholder bounds are a
    frame, not a clip. `normAutofit` is PowerPoint's own "shrink text on
    overflow", so the deck stays readable instead of overlapping, and it
    is the template's geometry doing the deciding rather than the author
    guessing how much will fit.
    """
    from pptx.enum.text import MSO_AUTO_SIZE

    try:
        placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        placeholder.text_frame.word_wrap = True
    except (AttributeError, ValueError):  # pragma: no cover - picture ph
        pass


def _fill_region(
    slide, placeholder, region: Region, base_dir: Path, build_dir: Path,
    code_typeface: str = DEFAULT_CODE_TYPEFACE,
) -> None:
    if region.type in ("text", "p", "ul", "ol"):
        _shrink_on_overflow(placeholder)
    if region.type == "text":
        # Plain string region: emphasis-aware, keeps the layout's own
        # paragraph style (titles and heads are unbulleted by design).
        write_runs(
            placeholder.text_frame.paragraphs[0],
            region.text or "",
            code_typeface=code_typeface,
        )
    elif region.type == "p":
        _fill_plain(placeholder, region, suppress_bullets=True, code_typeface=code_typeface)
    elif region.type in ("ul", "ol"):
        _fill_list(
            placeholder, region, numbered=(region.type == "ol"),
            code_typeface=code_typeface,
        )
    elif region.type == "image":
        _fill_image(
            slide, placeholder, _asset_path(base_dir, region.src),
            region.fit or DEFAULT_IMAGE_FIT,
        )
    elif region.type == "mermaid":
        image = resolve_mermaid(_asset_path(base_dir, region.src), build_dir)
        _fill_image(slide, placeholder, image, region.fit or DEFAULT_IMAGE_FIT)
    elif region.type == "video":
        _fill_video(slide, placeholder, region, base_dir, build_dir)
    else:  # pragma: no cover - parser guarantees the type set
        raise RenderError(f"unknown region type {region.type!r}")


def _render_slide(
    prs: Presentation,
    layouts_by_name: dict,
    binding: LayoutBinding,
    slide_src: Slide,
    session_id: str,
    base_dir: Path,
    build_dir: Path,
    code_typeface: str = DEFAULT_CODE_TYPEFACE,
    bindings: dict | None = None,
):
    layout = layouts_by_name[binding.layout_name]
    slide = prs.slides.add_slide(layout)

    placeholders_by_idx = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    for region in slide_src.regions:
        if region.name not in binding.regions:
            # Where that region *does* live. A region name is declared per
            # layout, so `full` working on one slide and not on another is
            # not a contradiction -- but the message said only that this
            # layout has not got it, which reads as the name being wrong
            # when the layout is the thing that is wrong.
            elsewhere = sorted(
                key
                for key, other in (bindings or {}).items()
                if region.name in getattr(other, "regions", {})
            )
            hint = (
                f" {region.name!r} is declared by {', '.join(repr(k) for k in elsewhere)}"
                f" — this slide may want one of those layouts instead."
                if elsewhere
                else f" No layout in this library declares {region.name!r}."
            )
            raise RenderError(
                f"slide {slide_src.id!r}: region {region.name!r} is not declared for "
                f"layout {binding.key!r}; expected one of {sorted(binding.regions)}."
                + hint
            )
        idx = binding.regions[region.name]
        placeholder = placeholders_by_idx.get(idx)
        if placeholder is None:  # pragma: no cover - layout validation catches this
            raise RenderError(
                f"slide {slide_src.id!r}: placeholder idx {idx} missing on new slide"
            )
        _fill_region(slide, placeholder, region, base_dir, build_dir, code_typeface)

    if slide_src.notes:
        slide.notes_slide.notes_text_frame.text = slide_src.notes

    creation_id = derive_creation_id(session_id, slide_src.id)
    inject_creation_id(slide, creation_id)
    return slide, creation_id


def render_session(
    session_path: Path,
    template_path: Path,
    layout_map_path: Path,
    out_dir: Path,
    attribution_path: Path | None = None,
    outcomes_path: Path | None = None,
    skip_bad_slides: bool = False,
) -> dict:
    """Render one session. Returns a summary dict with output paths.

    Attribution config auto-detects at the library root
    (<sessions dir>/../attribution.yaml) unless a path is given.
    """
    session = parse_session(session_path)
    bindings = load_layout_map(layout_map_path)
    code_typeface = load_style(layout_map_path).get("code_typeface", DEFAULT_CODE_TYPEFACE)

    if attribution_path is None:
        candidate = session_path.parent.parent / "attribution.yaml"
        attribution_path = candidate if candidate.exists() else None
    attribution = load_attribution(attribution_path) if attribution_path else None

    # The outcome catalogue, found the same way. Optional: a library
    # without one derives nothing and renders exactly as it did before.
    if outcomes_path is None:
        candidate = session_path.parent.parent / OUTCOMES_FILE
        outcomes_path = candidate if candidate.exists() else None
    catalogue = load_outcomes(outcomes_path) if outcomes_path else {}

    # A slide with a role is filled from the library rather than typed:
    # the session title and its outcomes are already written down, and a
    # second copy on a slide is a second thing to keep in step. Done here,
    # before anything reads the regions, so the deck, the index and the
    # catalog all see the same filled slide.
    block_meta_path = session_path.parent / "block.yaml"
    block_meta = {}
    if block_meta_path.is_file():
        try:
            block_meta = yamlio.safe_load(block_meta_path.read_text(encoding="utf-8")) or {}
        except yaml.YAMLError:
            block_meta = {}
    fill_roles(
        session,
        block_meta,
        catalogue,
        {key: list(binding.regions) for key, binding in bindings.items()},
    )

    prs = Presentation(str(template_path))
    validate_against_template(bindings, prs, template_path)

    layouts_by_name = {
        layout.name: layout for master in prs.slide_masters for layout in master.slide_layouts
    }

    base_dir = session_path.parent
    build_dir = out_dir / "build"
    out_dir.mkdir(parents=True, exist_ok=True)

    session_id = session.session_id
    pptx_name = f"session-{session_id}.pptx"

    slide_id_map: dict[str, str] = {}
    index_entries: list[dict] = []

    skipped: list[dict] = []
    for slide_src in session.slides:
        binding = bindings.get(slide_src.layout)
        try:
            if binding is None:
                raise RenderError(
                    f"slide {slide_src.id!r}: layout {slide_src.layout!r} not in layout map; "
                    f"available: {sorted(bindings)}"
                )
            slide, creation_id = _render_slide(
                prs, layouts_by_name, binding, slide_src, session_id, base_dir,
                build_dir, code_typeface, bindings,
            )
        except RenderError as err:
            # Strict by default: a build that quietly drops a slide is a
            # deck that is wrong in a way nobody notices until the room.
            if not skip_bad_slides:
                raise
            # Asked to carry on: the slide is left out and said so. One
            # slide nobody can fix today should not stop the other forty
            # being used, which is what a hard failure means for anyone
            # assembling a deck from somebody else's library.
            skipped.append({"slideId": slide_src.id, "message": str(err)})
            continue

        dc = session.frontmatter.get("dc") or {}
        meta = {
            "creator": dc.get("creator"),
            "license": dc.get("license"),
            "session": session_id,
            "slide": slide_src.id,
            "version": session.frontmatter.get("version"),
        }
        inject_provenance(slide, meta)
        if attribution:
            stamp_slide(slide, prs, attribution, build_dir, meta)
            append_notes_line(slide, attribution["text"], session_id, slide_src.id)

        source_ref = f"{slide.slide_id}#{creation_id}"
        slide_id_map[slide_src.id] = source_ref

        title_region = next((r for r in slide_src.regions if r.name == "title"), None)
        plain_title = plain_text(title_region.text) if title_region and title_region.text else None
        index_entries.append(
            {
                "slideId": slide_src.id,
                "sessionId": session_id,
                "sourceRef": source_ref,
                "sourcePptx": pptx_name,
                "title": plain_title,
                # Derived: what the slide declares, plus what its outcomes
                # develop. Everything downstream -- catalog.json, the pane,
                # tracking, credential checking -- reads this one field and
                # needed no change. `outcomes` rides alongside so the
                # provenance stays recoverable.
                "develops": derive_develops(
                    slide_src.develops, slide_src.outcomes, catalogue
                ),
                "outcomes": slide_src.outcomes,
                "dok": slide_src.dok,
                # Carried so the pane can search narration, not just titles.
                "notes": slide_src.notes,
            }
        )

    pptx_path = out_dir / pptx_name
    prs.save(str(pptx_path))
    normalize_zip(pptx_path)

    map_path = out_dir / "slide-id-map.json"
    map_path.write_text(
        json.dumps(slide_id_map, indent=2, sort_keys=True) + "\n", encoding="utf-8"
    )

    index_path = out_dir / "index.json"
    index_doc = {
        "session": {
            key: value
            for key, value in session.frontmatter.items()
        },
        "slides": index_entries,
    }
    index_path.write_text(json.dumps(index_doc, indent=2) + "\n", encoding="utf-8")

    return {
        "pptx": pptx_path,
        "slide_id_map": map_path,
        "index": index_path,
        "slide_count": len(session.slides) - len(skipped),
        # Empty unless skip_bad_slides was asked for and something used
        # it. Carried so a caller can say which slides are missing and
        # why, rather than a deck quietly being shorter than the file.
        "skipped": skipped,
    }
