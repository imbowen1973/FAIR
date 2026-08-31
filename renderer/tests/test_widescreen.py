"""Widening a template, without squaring the scale.

The first attempt at this scaled through python-pptx's wrappers, which
reach the same offset more than once: 1.333 x 1.333 is 1.778, and
placeholders ended up hanging off the right-hand edge. So the test that
matters is not "did it get wider" but "is every placeholder exactly where
proportion says it should be, and still on the slide".
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from edufair_renderer.widescreen import STANDARD_EMU, WIDESCREEN_EMU, widen

EXAMPLES = Path(__file__).resolve().parents[2] / "examples"


def _narrow(tmp_path: Path) -> Path:
    """The example template, put back to 4:3, to widen again."""
    src = tmp_path / "template.pptx"
    shutil.copy(EXAMPLES / "template.pptx", src)
    prs = Presentation(str(src))
    if prs.slide_width == WIDESCREEN_EMU:
        widen(src, width=STANDARD_EMU)
    return src


def _placeholders(path: Path):
    prs = Presentation(str(path))
    out = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for ph in layout.placeholders:
                out[(layout.name, ph.placeholder_format.idx)] = (
                    ph.left, ph.top, ph.width, ph.height
                )
    return out, prs.slide_width, prs.slide_height


def test_the_slide_gets_wider_and_no_taller(tmp_path):
    src = _narrow(tmp_path)
    before, was_w, was_h = _placeholders(src)
    assert was_w == STANDARD_EMU

    result = widen(src)
    after, now_w, now_h = _placeholders(src)

    assert result["changed"] is True
    assert now_w == WIDESCREEN_EMU
    assert now_h == was_h, "the height does not change, so nothing vertical should"


def test_every_placeholder_keeps_its_proportions(tmp_path):
    # The squared-scale bug, caught by arithmetic rather than by eye.
    src = _narrow(tmp_path)
    before, was_w, _ = _placeholders(src)
    widen(src)
    after, now_w, _ = _placeholders(src)

    factor = now_w / was_w
    assert before, "the example template has no placeholders to check"
    for key, (left, top, width, height) in before.items():
        new_left, new_top, new_width, new_height = after[key]
        assert new_left == pytest.approx(left * factor, abs=2), key
        assert new_width == pytest.approx(width * factor, abs=2), key
        # Vertical is untouched: the slide is no taller.
        assert new_top == top, key
        assert new_height == height, key


def test_nothing_hangs_off_the_edge(tmp_path):
    # What the squared scale actually looked like from the outside.
    src = _narrow(tmp_path)
    widen(src)
    after, width, height = _placeholders(src)
    for key, (left, top, w, h) in after.items():
        assert left + w <= width + 2, f"{key} runs past the right-hand edge"
        assert top + h <= height + 2, f"{key} runs past the bottom"


def test_widening_twice_changes_nothing_the_second_time(tmp_path):
    # Running it again is how somebody checks whether they ran it, and it
    # must not scale a template that is already wide.
    src = _narrow(tmp_path)
    widen(src)
    once, _, _ = _placeholders(src)
    result = widen(src)
    twice, width, _ = _placeholders(src)
    assert result["changed"] is False
    assert width == WIDESCREEN_EMU
    assert once == twice


def test_the_deck_still_renders_after_widening(tmp_path):
    # A template that is wider but broken is not an improvement.
    import textwrap
    import yaml
    from edufair_renderer.corpus import build_corpus

    root = tmp_path / "lib"
    root.mkdir()
    shutil.copy(EXAMPLES / "layout-map.yaml", root / "layout-map.yaml")
    src = _narrow(root)
    widen(src)

    block = root / "blocks" / "01-b"
    block.mkdir(parents=True)
    (block / "block.yaml").write_text(
        yaml.safe_dump({"title": "B", "duration_minutes": 45, "resources": []}),
        encoding="utf-8",
    )
    (block / "slides.md").write_text(
        textwrap.dedent(
            """\
            ---
            session: '01'
            title: B
            version: 1.0.0
            ---

            --- slide
            id: s-01
            layout: Full
            title: Wide
            full: Still renders
            ---
            """
        ),
        encoding="utf-8",
    )
    (root / "course.yaml").write_text(
        yaml.safe_dump({"id": "c", "title": "C", "structure": [{"block": "01-b"}]}),
        encoding="utf-8",
    )

    summary = build_corpus(
        library=root,
        template=src,
        layout_map=root / "layout-map.yaml",
        out_dir=tmp_path / "out",
        warn=lambda m: None,
    )
    assert summary["slides"] == 1
