"""Acceptance tests for spec A.6, plus the creationId contract.

Criterion numbering follows docs/md-to-powerpoint-pipeline.md section A.6.
"""

from __future__ import annotations

import json
import re
import zipfile

import yaml
from pathlib import Path

import pytest
from lxml import etree

import shutil

from pptx import Presentation

from edufair_renderer.corpus import main as corpus_main
from edufair_renderer.render import render_session

REPO = Path(__file__).resolve().parents[2]
EXAMPLES = REPO / "examples"

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS = {"p": P_NS, "a": A_NS, "p14": P14_NS}


@pytest.fixture(scope="module")
def rendered(tmp_path_factory):
    out_dir = tmp_path_factory.mktemp("out")
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-01.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=out_dir,
    )
    return result


def _slide_xml(pptx_path: Path, n: int) -> etree._Element:
    with zipfile.ZipFile(pptx_path) as z:
        return etree.fromstring(z.read(f"ppt/slides/slide{n}.xml"))


def _slide_layout_name(pptx_path: Path, n: int) -> str:
    with zipfile.ZipFile(pptx_path) as z:
        rels = etree.fromstring(z.read(f"ppt/slides/_rels/slide{n}.xml.rels"))
        layout_target = None
        for rel in rels:
            if rel.get("Type").endswith("/slideLayout"):
                layout_target = rel.get("Target").replace("../", "ppt/")
        assert layout_target, f"slide{n} has no slideLayout relationship"
        layout = etree.fromstring(z.read(layout_target))
        return layout.find(f"{{{P_NS}}}cSld").get("name")


def test_a6_1_slide_references_split_layout(rendered):
    """Slide s01-03 (deck slide 3) must reference the Split layout."""
    assert _slide_layout_name(rendered["pptx"], 3) == "Split"


def test_a6_2_placeholder_indices_present(rendered):
    """Slide XML must carry <p:ph> references with idx 0, 1, 2.

    (idx attribute is omitted for idx 0 title placeholders per OOXML —
    a <p:ph type="title"> with no idx attribute IS idx 0.)
    """
    root = _slide_xml(rendered["pptx"], 3)
    phs = root.findall(f".//{{{P_NS}}}ph")
    indices = set()
    for ph in phs:
        idx = ph.get("idx")
        indices.add(int(idx) if idx is not None else 0)
    assert indices == {0, 1, 2}, f"expected placeholder idx {{0,1,2}}, got {indices}"


def test_a6_3_content_in_correct_placeholders(rendered):
    """Left content in idx 1, right content in idx 2 — not swapped or merged."""
    root = _slide_xml(rendered["pptx"], 3)
    text_by_idx = {}
    for sp in root.findall(f".//{{{P_NS}}}sp"):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None:
            continue
        idx = int(ph.get("idx") or 0)
        texts = [t.text or "" for t in sp.findall(f".//{{{A_NS}}}t")]
        text_by_idx[idx] = texts
    assert "Temperature logging" in text_by_idx[1]
    assert "Sensor thresholds" in text_by_idx[1]
    assert "Receive shipment" in text_by_idx[2]
    assert "Temperature logging" not in text_by_idx[2]
    assert "Receive shipment" not in text_by_idx[1]
    assert text_by_idx[0] == ["Cold chain integrity"]


def test_a6_3b_nesting_levels(rendered):
    """Nested ul items must carry paragraph level attributes."""
    root = _slide_xml(rendered["pptx"], 3)
    for sp in root.findall(f".//{{{P_NS}}}sp"):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None or int(ph.get("idx") or 0) != 1:
            continue
        levels = {}
        for p in sp.findall(f".//{{{A_NS}}}p"):
            texts = [t.text or "" for t in p.findall(f".//{{{A_NS}}}t")]
            pPr = p.find(f"{{{A_NS}}}pPr")
            lvl = int(pPr.get("lvl")) if pPr is not None and pPr.get("lvl") else 0
            if texts:
                levels["".join(texts)] = lvl
        assert levels["Temperature logging"] == 0
        assert levels["Sensor thresholds"] == 1
        assert levels["Manual inspection"] == 1
        assert levels["Audit trail"] == 0
        return
    pytest.fail("left placeholder (idx 1) not found")


def test_a6_4_text_is_editable_runs(rendered):
    """Text must be <a:t> runs inside <p:sp> shapes, not rasterised images."""
    root = _slide_xml(rendered["pptx"], 3)
    runs = root.findall(f".//{{{P_NS}}}sp//{{{A_NS}}}t")
    assert len(runs) >= 7
    pics = root.findall(f".//{{{P_NS}}}pic")
    assert pics == [], "text slide must not contain rasterised pictures"


def test_a6_5_deterministic_output(rendered, tmp_path):
    """Two runs of the same input must produce byte-identical files."""
    out2 = tmp_path / "run2"
    result2 = render_session(
        session_path=EXAMPLES / "sessions" / "session-01.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=out2,
    )
    assert rendered["pptx"].read_bytes() == result2["pptx"].read_bytes()
    assert (
        rendered["slide_id_map"].read_text() == result2["slide_id_map"].read_text()
    )
    assert rendered["index"].read_text() == result2["index"].read_text()


def test_creation_id_present_and_mapped(rendered):
    """Every slide carries p14:creationId, and slide-id-map.json matches the XML."""
    slide_id_map = json.loads(rendered["slide_id_map"].read_text())
    assert set(slide_id_map) == {"s01-01", "s01-02", "s01-03", "s01-04", "s01-05"}

    with zipfile.ZipFile(rendered["pptx"]) as z:
        # presentation.xml gives authored order of slide ids
        pres = etree.fromstring(z.read("ppt/presentation.xml"))
        sld_ids = [
            int(el.get("id"))
            for el in pres.findall(f".//{{{P_NS}}}sldIdLst/{{{P_NS}}}sldId")
        ]
        creation_ids = []
        for n in range(1, 6):
            root = etree.fromstring(z.read(f"ppt/slides/slide{n}.xml"))
            cid = root.find(f".//{{{P14_NS}}}creationId")
            assert cid is not None, f"slide{n} missing p14:creationId"
            creation_ids.append(int(cid.get("val")))

    expected_refs = {f"{sid}#{cid}" for sid, cid in zip(sld_ids, creation_ids)}
    assert set(slide_id_map.values()) == expected_refs
    for ref in slide_id_map.values():
        assert re.fullmatch(r"\d+#\d+", ref), f"bad sourceRef format: {ref}"


def test_creation_id_is_stable():
    from edufair_renderer.creation_id import derive_creation_id

    a = derive_creation_id("01", "s01-03")
    b = derive_creation_id("01", "s01-03")
    assert a == b
    assert 1 <= a <= 0x7FFFFFFF
    assert derive_creation_id("01", "s01-04") != a


def test_index_json_contract(rendered):
    """index.json carries the semantic fields the assembler and graph need."""
    doc = json.loads(rendered["index"].read_text())
    assert doc["session"]["session"] == "01"
    slides = {s["slideId"]: s for s in doc["slides"]}
    entry = slides["s01-03"]
    assert entry["sessionId"] == "01"
    assert entry["sourcePptx"] == "session-01.pptx"
    assert entry["title"] == "Cold chain integrity"
    assert entry["develops"] == ["C1", "C4"]
    assert entry["dok"] == 2
    assert re.fullmatch(r"\d+#\d+", entry["sourceRef"])
    # Notes ride in the index so the pane can search narration, not just
    # titles; slides without notes carry an explicit null, never a gap.
    assert "integrity pillars" in entry["notes"]
    assert all("notes" in s for s in doc["slides"])
    assert any(s["notes"] is None for s in doc["slides"])


def test_notes_written(rendered):
    with zipfile.ZipFile(rendered["pptx"]) as z:
        names = [n for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide")]
        assert names, "no notes slides emitted"
        found = any(
            b"integrity pillars" in z.read(n) for n in names
        )
    assert found, "speaker notes text not found in any notes slide"


def test_picture_placeholder_keeps_binding(rendered):
    """The Picture layout's image binds into the picture placeholder —
    the <p:pic> itself carries <p:ph idx=1>, no free-positioning."""
    root = _slide_xml(rendered["pptx"], 4)
    pics = root.findall(f".//{{{P_NS}}}pic")
    assert len(pics) == 1
    ph = pics[0].find(f".//{{{P_NS}}}ph")
    assert ph is not None, "picture must retain its placeholder reference"
    assert int(ph.get("idx") or 0) == 1


def test_caption_plain_text_suppresses_bullets(rendered):
    """The p-type caption renders paragraphs with buNone and bold emphasis."""
    root = _slide_xml(rendered["pptx"], 4)
    for sp in root.findall(f".//{{{P_NS}}}sp"):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None or int(ph.get("idx") or 0) != 2:
            continue
        paragraphs = sp.findall(f".//{{{A_NS}}}p")
        assert len(paragraphs) == 2
        for p in paragraphs:
            assert p.find(f"{{{A_NS}}}pPr/{{{A_NS}}}buNone") is not None
        bold_runs = [
            r for r in sp.findall(f".//{{{A_NS}}}r")
            if r.find(f"{{{A_NS}}}rPr") is not None
            and r.find(f"{{{A_NS}}}rPr").get("b") == "1"
        ]
        assert any((r.find(f"{{{A_NS}}}t").text or "") == "red line" for r in bold_runs)
        return
    pytest.fail("caption placeholder (idx 2) not found on Picture slide")


def test_theme_color_and_emphasis_on_lists(rendered):
    """Region/item colours resolve to schemeClr runs; ** renders as bold."""
    root = _slide_xml(rendered["pptx"], 5)
    for sp in root.findall(f".//{{{P_NS}}}sp"):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None or int(ph.get("idx") or 0) != 1:
            continue
        colors_by_text = {}
        bold_texts = set()
        for r in sp.findall(f".//{{{A_NS}}}r"):
            text = r.find(f"{{{A_NS}}}t").text or ""
            rPr = r.find(f"{{{A_NS}}}rPr")
            if rPr is None:
                continue
            scheme = rPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
            if scheme is not None:
                colors_by_text[text] = scheme.get("val")
            if rPr.get("b") == "1":
                bold_texts.add(text)
        # region default accent1; item override accent2
        assert colors_by_text["In range"] == "accent1"
        assert colors_by_text["Excursion"] == "accent2"
        assert "In range" in bold_texts and "Excursion" in bold_texts
        return
    pytest.fail("left placeholder (idx 1) not found on slide 5")


def test_comparison_layout_binds_heads_and_bodies(tmp_path):
    """Comparison slide binds all five placeholder indices."""
    out = tmp_path / "cmp-out"
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-02.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=out,
    )
    assert _slide_layout_name(result["pptx"], 3) == "Comparison"
    root = _slide_xml(result["pptx"], 3)
    text_by_idx = {}
    for sp in root.findall(f".//{{{P_NS}}}sp"):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None:
            continue
        idx = int(ph.get("idx") or 0)
        text_by_idx[idx] = [t.text or "" for t in sp.findall(f".//{{{A_NS}}}t")]
    assert set(text_by_idx) == {0, 1, 2, 3, 4}
    assert text_by_idx[1] == ["Paper"]
    assert text_by_idx[3] == ["Digital"]
    assert "Manual transcription errors" in text_by_idx[2]
    assert "Automatic logger export" in text_by_idx[4]


def test_cards_layout_binds_four_tabbed_cards(tmp_path):
    """Cards slide binds title + 4 tab heads (idx 1-4) + 4 bodies (idx 5-8),
    and the layout styles tabs with per-card accent fills."""
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-03.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=tmp_path / "out",
    )
    assert _slide_layout_name(result["pptx"], 4) == "Cards"
    root = _slide_xml(result["pptx"], 4)
    text_by_idx = {}
    for sp in root.findall(f".//{{{P_NS}}}sp"):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None:
            continue
        idx = int(ph.get("idx") or 0)
        text_by_idx[idx] = [t.text or "" for t in sp.findall(f".//{{{A_NS}}}t")]
    assert set(text_by_idx) >= {0, 1, 2, 3, 4, 5, 6, 7, 8}
    assert text_by_idx[1] == ["Detect"]
    assert text_by_idx[4] == ["Report"]
    assert "Quarantine the stock" in text_by_idx[6]
    assert "File the excursion report" in text_by_idx[8]

    # Layout: tab placeholders carry rounded-top geometry and accent fills.
    with zipfile.ZipFile(EXAMPLES / "template.pptx") as z:
        cards_layout = None
        for name in z.namelist():
            if name.startswith("ppt/slideLayouts/slideLayout") and name.endswith(".xml"):
                xml = etree.fromstring(z.read(name))
                if xml.find(f"{{{P_NS}}}cSld").get("name") == "Cards":
                    cards_layout = xml
                    break
    assert cards_layout is not None
    fills = {}
    for sp in cards_layout.findall(f".//{{{P_NS}}}sp"):
        ph = sp.find(f".//{{{P_NS}}}ph")
        if ph is None or ph.get("idx") is None:
            continue
        idx = int(ph.get("idx"))
        if 1 <= idx <= 4:
            geom = sp.find(f".//{{{A_NS}}}prstGeom")
            assert geom is not None and geom.get("prst") == "round2SameRect"
            scheme = sp.find(f".//{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
            fills[idx] = scheme.get("val")
    assert fills == {1: "accent1", 2: "accent2", 3: "accent3", 4: "accent4"}


def test_video_region_renders_linked_poster(tmp_path):
    """Video regions render a poster picture hyperlinked to the hosted URL;
    no video binary is involved anywhere."""
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-02.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=tmp_path / "out",
    )
    with zipfile.ZipFile(result["pptx"]) as z:
        root = etree.fromstring(z.read("ppt/slides/slide4.xml"))
        pics = root.findall(f".//{{{P_NS}}}pic")
        assert len(pics) == 1
        hlink = pics[0].find(f".//{{{A_NS}}}hlinkClick")
        assert hlink is not None, "poster picture must carry a hyperlink"
        rid = hlink.get(f"{{{R_NS}}}id")
        rels = etree.fromstring(z.read("ppt/slides/_rels/slide4.xml.rels"))
        targets = {r.get("Id"): r.get("Target") for r in rels}
        assert targets[rid] == "https://youtu.be/fair-demo-audit-pack"


def test_asset_policy_gate(tmp_path):
    import sys

    sys.path.insert(0, str(REPO / "scripts"))
    from check_assets import check_assets
    from PIL import Image

    # clean directory passes
    assert check_assets(EXAMPLES / "sessions") == []

    bad = tmp_path / "assets"
    bad.mkdir()
    (bad / "clip.mp4").write_bytes(b"\x00" * 10)
    Image.new("RGB", (3000, 200), "white").save(bad / "huge.png")
    errors = check_assets(tmp_path)
    assert any("clip.mp4" in e and "video" in e for e in errors)
    assert any("huge.png" in e and "2200" in e for e in errors)


FAIR_NS = "urn:fair:attribution:1"


def test_attribution_stamp_locked_and_present(rendered):
    """Every slide carries the visible CC-BY stamp as an unselectable shape."""
    for n in range(1, 6):
        root = _slide_xml(rendered["pptx"], n)
        stamped = False
        for sp in root.findall(f".//{{{P_NS}}}sp"):
            texts = "".join(t.text or "" for t in sp.findall(f".//{{{A_NS}}}t"))
            if "CC-BY 4.0" in texts:
                locks = sp.find(f".//{{{A_NS}}}spLocks")
                assert locks is not None, f"slide{n}: stamp has no locks"
                for attr in ("noSelect", "noMove", "noResize", "noTextEdit"):
                    assert locks.get(attr) == "1", f"slide{n}: missing {attr}"
                stamped = True
        assert stamped, f"slide{n}: no visible attribution stamp"


def test_attribution_provenance_extension(rendered):
    """Machine-readable provenance rides in every slide's XML."""
    root = _slide_xml(rendered["pptx"], 3)
    node = root.find(f".//{{{FAIR_NS}}}attribution")
    assert node is not None
    assert node.get("creator") == "FAIR Consortium"
    assert node.get("license") == "CC-BY-4.0"
    assert node.get("session") == "01"
    assert node.get("slide") == "s01-03"


def test_attribution_notes_line(rendered):
    with zipfile.ZipFile(rendered["pptx"]) as z:
        notes = b"".join(
            z.read(n) for n in z.namelist() if n.startswith("ppt/notesSlides/notesSlide")
        )
    assert b"Source: Example Grant" in notes
    assert b"01/s01-03" in notes


def test_audit_reads_provenance(rendered):
    from edufair_renderer.audit import audit_deck

    entries = audit_deck(rendered["pptx"])
    assert len(entries) == 5
    by_slide = {e["provenance"]["slide"]: e for e in entries if e["provenance"]}
    assert by_slide["s01-03"]["provenance"]["license"] == "CC-BY-4.0"
    assert all(e["creationId"] for e in entries)


def test_attribution_logo_stamped_with_metadata(tmp_path):
    """With a logo configured, every slide gets a locked picture whose PNG
    carries the attribution in its own metadata."""
    from PIL import Image

    lib = tmp_path / "lib"
    (lib / "branding").mkdir(parents=True)
    Image.new("RGBA", (120, 120), (44, 110, 158, 255)).save(lib / "branding" / "logo.png")
    # Explicit utf-8: the loader reads utf-8, and on Windows the default
    # locale encoding is cp1252, which mangles the middle dot on write.
    (lib / "attribution.yaml").write_text(
        'text: "CC-BY 4.0 · FAIR Consortium"\nlogo: branding/logo.png\n',
        encoding="utf-8",
    )
    out = tmp_path / "out"
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-02.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=out,
        attribution_path=lib / "attribution.yaml",
    )
    root = _slide_xml(result["pptx"], 1)
    pics = root.findall(f".//{{{P_NS}}}pic")
    assert pics, "logo picture missing"
    lock_found = any(
        pic.find(f".//{{{A_NS}}}picLocks") is not None
        and pic.find(f".//{{{A_NS}}}picLocks").get("noSelect") == "1"
        for pic in pics
    )
    assert lock_found, "logo picture is not locked"

    stamped = Image.open(out / "build" / "attribution-logo.png")
    assert stamped.info.get("fair:license") or stamped.info.get("fair:text")


def test_master_bullet_levels_colored():
    """Template master colour-codes list levels 1-3 (accent1/2/3)."""
    with zipfile.ZipFile(EXAMPLES / "template.pptx") as z:
        master = etree.fromstring(z.read("ppt/slideMasters/slideMaster1.xml"))
    body = master.find(f"{{{P_NS}}}txStyles/{{{P_NS}}}bodyStyle")
    for level, expected in ((1, "accent1"), (2, "accent2"), (3, "accent3")):
        lvl = body.find(f"{{{A_NS}}}lvl{level}pPr")
        scheme = lvl.find(f"{{{A_NS}}}buClr/{{{A_NS}}}schemeClr")
        assert scheme is not None and scheme.get("val") == expected


def test_layout_map_validation_fails_loudly(tmp_path):
    """A stale placeholder index must abort the build, not misbind."""
    from edufair_renderer.layoutmap import LayoutMapError

    bad_map = tmp_path / "bad-map.yaml"
    bad_map.write_text(
        "Split:\n  layout_name: 'Split'\n  regions:\n    title: 0\n    left: 1\n    right: 9\n"
    )
    with pytest.raises(LayoutMapError, match="idx 9"):
        render_session(
            session_path=EXAMPLES / "sessions" / "session-01.md",
            template_path=EXAMPLES / "template.pptx",
            layout_map_path=bad_map,
            out_dir=tmp_path / "out",
        )


# --- nested delivery structure (credentials modules/days/blocks) ----------


def _corpus(tmp_path, credential: dict):
    """Build a one-session corpus with the given credential definition."""
    import shutil

    import yaml

    from edufair_renderer.corpus import build_corpus

    sessions = tmp_path / "sessions"
    if not sessions.exists():
        # copytree brings assets/ along; session-01 references an image.
        shutil.copytree(EXAMPLES / "sessions", sessions)
        for stale in sessions.glob("session-0[23].md"):
            stale.unlink()
    creds = tmp_path / "credentials"
    creds.mkdir(exist_ok=True)
    (creds / "c.yaml").write_text(yaml.safe_dump(credential), encoding="utf-8")
    return build_corpus(
        sessions_dir=sessions,
        template=EXAMPLES / "template.pptx",
        layout_map=EXAMPLES / "layout-map.yaml",
        out_dir=tmp_path / "out",
        credentials_dir=creds,
        warn=lambda msg: None,
    )


def test_nested_credential_validates(tmp_path):
    """modules -> days -> blocks -> sessions is accepted and passed through."""
    summary = _corpus(
        tmp_path,
        {
            "id": "c",
            "title": "Course",
            "modules": [
                {
                    "title": "M1",
                    "days": [
                        {"title": "Day 1", "blocks": [{"title": "AM", "sessions": ["01"]}]}
                    ],
                }
            ],
        },
    )
    assert summary["credentials"] == 1
    catalog = json.loads((tmp_path / "out" / "catalog.json").read_text())
    block = catalog["credentials"][0]["modules"][0]["days"][0]["blocks"][0]
    assert block["sessions"] == ["01"]


def test_flat_credential_still_validates(tmp_path):
    """The pre-existing flat modules[].sessions[] shape must not break."""
    summary = _corpus(
        tmp_path,
        {"id": "c", "title": "Course", "modules": [{"title": "M1", "sessions": ["01"]}]},
    )
    assert summary["credentials"] == 1


def test_unknown_session_in_nested_block_raises(tmp_path):
    from edufair_renderer.corpus import CorpusError

    with pytest.raises(CorpusError, match="unknown session"):
        _corpus(
            tmp_path,
            {
                "id": "c",
                "modules": [
                    {
                        "title": "M1",
                        "days": [
                            {"title": "D", "blocks": [{"title": "AM", "sessions": ["99"]}]}
                        ],
                    }
                ],
            },
        )


def test_container_with_both_sessions_and_days_raises(tmp_path):
    from edufair_renderer.corpus import CorpusError

    with pytest.raises(CorpusError, match="one or the other"):
        _corpus(
            tmp_path,
            {
                "id": "c",
                "modules": [
                    {
                        "title": "M1",
                        "sessions": ["01"],
                        "days": [{"title": "D", "sessions": ["01"]}],
                    }
                ],
            },
        )


def test_attribution_opacity_and_scale_reach_the_xml(tmp_path):
    """Semi-transparency is a render concern: alpha lands in the slide XML."""
    from PIL import Image

    lib = tmp_path / "lib"
    (lib / "branding").mkdir(parents=True)
    Image.new("RGBA", (240, 80), (44, 110, 158, 255)).save(lib / "branding" / "logo.png")
    (lib / "attribution.yaml").write_text(
        'text: "Funded by the European Union · Agreement No: 101183898"\n'
        "logo: branding/logo.png\n"
        "corner: bottom-right\n"
        "opacity: 0.45\n"
        "scale: 0.9\n"
        # Fading a mark is opt-in: logo_opacity defaults to fully opaque
        # so an emblem is never altered by accident.
        "logo_opacity: 0.45\n",
        encoding="utf-8",
    )
    out = tmp_path / "out"
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-02.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=out,
        attribution_path=lib / "attribution.yaml",
    )
    with zipfile.ZipFile(result["pptx"]) as z:
        xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
    # 0.45 -> 45000 in DrawingML's 0-100000 alpha scale.
    assert 'alphaModFix amt="45000"' in xml, "logo is not faded"
    assert '<a:alpha val="45000"/>' in xml, "stamp text is not faded"


def test_attribution_rejects_out_of_range_opacity(tmp_path):
    from edufair_renderer.attribution import AttributionError, load_attribution

    path = tmp_path / "attribution.yaml"
    path.write_text('text: "x"\nopacity: 1.5\n', encoding="utf-8")
    with pytest.raises(AttributionError, match="opacity"):
        load_attribution(path)


def test_attribution_defaults_to_opaque(tmp_path):
    from edufair_renderer.attribution import load_attribution

    path = tmp_path / "attribution.yaml"
    path.write_text('text: "x"\n', encoding="utf-8")
    config = load_attribution(path)
    assert config["opacity"] == 1.0 and config["scale"] == 1.0


def test_catalog_carries_attribution_and_copies_the_logo(tmp_path):
    """The funder credit travels with the content, not with the tool."""
    import shutil

    from PIL import Image

    from edufair_renderer.corpus import build_corpus

    lib = tmp_path / "lib"
    shutil.copytree(EXAMPLES / "sessions", lib / "sessions")
    for stale in (lib / "sessions").glob("session-0[23].md"):
        stale.unlink()
    (lib / "branding").mkdir()
    Image.new("RGBA", (200, 60), (0, 51, 153, 255)).save(lib / "branding" / "eu.png")
    (lib / "attribution.yaml").write_text(
        'text: "Funded by the European Union · Agreement No: 101183898"\n'
        "logo: branding/eu.png\nopacity: 0.45\n",
        encoding="utf-8",
    )

    out = tmp_path / "out"
    # attribution.yaml sits beside sessions/, so it is found without asking.
    build_corpus(
        sessions_dir=lib / "sessions",
        template=EXAMPLES / "template.pptx",
        layout_map=EXAMPLES / "layout-map.yaml",
        out_dir=out,
        warn=lambda msg: None,
    )
    catalog = json.loads((out / "catalog.json").read_text())
    credit = catalog["attribution"]
    assert "101183898" in credit["text"]
    assert credit["opacity"] == 0.45
    # Referenced relatively and actually present, so a consumer needs
    # nothing but the data directory.
    assert credit["logo"] == "attribution-logo.png"
    assert (out / "attribution-logo.png").exists()


def test_catalog_attribution_is_null_without_a_config(tmp_path):
    import shutil

    from edufair_renderer.corpus import build_corpus

    lib = tmp_path / "lib"
    shutil.copytree(EXAMPLES / "sessions", lib / "sessions")
    out = tmp_path / "out"
    build_corpus(
        sessions_dir=lib / "sessions",
        template=EXAMPLES / "template.pptx",
        layout_map=EXAMPLES / "layout-map.yaml",
        out_dir=out,
        warn=lambda msg: None,
    )
    catalog = json.loads((out / "catalog.json").read_text())
    assert catalog["attribution"] is None


def test_logo_sizing_and_opacity_are_independent_of_the_text(tmp_path):
    """EU rules: emblem >= 1 cm and unaltered, while the wording stays faded."""
    from PIL import Image

    lib = tmp_path / "lib"
    (lib / "branding").mkdir(parents=True)
    # 4.75:1, the real emblem's aspect ratio.
    Image.new("RGBA", (1130, 238), (0, 51, 153, 255)).save(lib / "branding" / "eu.png")
    (lib / "attribution.yaml").write_text(
        'text: "Co-funded by the European Union · Agreement No: 101183898"\n'
        "logo: branding/eu.png\nopacity: 0.45\nscale: 0.9\n"
        "logo_height_cm: 1.0\nlogo_opacity: 1.0\n",
        encoding="utf-8",
    )
    out = tmp_path / "out"
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-02.md",
        template_path=EXAMPLES / "template.pptx",
        layout_map_path=EXAMPLES / "layout-map.yaml",
        out_dir=out,
        attribution_path=lib / "attribution.yaml",
    )
    with zipfile.ZipFile(result["pptx"]) as z:
        xml = z.read("ppt/slides/slide1.xml").decode("utf-8")

    # The wording is faded; the emblem is not touched.
    assert '<a:alpha val="45000"/>' in xml, "stamp text should be faded"
    assert "alphaModFix" not in xml, "the EU emblem must not be altered"

    # 1 cm = 360000 EMU, within rounding.
    root = etree.fromstring(xml.encode("utf-8"))
    ext = root.find(f".//{{{P_NS}}}pic/{{{P_NS}}}spPr/{{{A_NS}}}xfrm/{{{A_NS}}}ext")
    assert ext is not None, "no picture on the slide"
    cy = int(ext.get("cy"))
    assert abs(cy - 360000) < 2000, f"emblem is {cy / 360000:.2f} cm high, not 1 cm"


# --- template profile / edufair-template ------------------------------------

STOCK = Path(__import__("pptx").__file__).parent / "templates" / "default.pptx"


def test_stock_powerpoint_template_conforms_to_the_core_profile():
    """The whole point: a plain PowerPoint template works unedited."""
    from edufair_renderer.template_inspect import CORE, inspect_template

    results = {r.key: r for r in inspect_template(STOCK)}
    for spec in CORE:
        result = results[spec.key]
        assert result.layout_name is not None, f"{spec.key} unmatched in a stock template"
        assert not result.missing, f"{spec.key} unbound regions: {result.missing}"


def test_generated_map_renders_real_content_on_a_stock_template(tmp_path):
    from edufair_renderer.layoutmap import load_layout_map
    from edufair_renderer.template_inspect import inspect_template, to_layout_map

    map_path = tmp_path / "layout-map.yaml"
    map_path.write_text(to_layout_map(inspect_template(STOCK)), encoding="utf-8")
    # It must survive the same loader the build uses, then actually bind.
    load_layout_map(map_path)
    result = render_session(
        session_path=EXAMPLES / "sessions" / "session-02.md",
        template_path=STOCK,
        layout_map_path=map_path,
        out_dir=tmp_path / "out",
    )
    assert result["pptx"].exists()


def test_inspector_reproduces_the_hand_written_map():
    """The generator agrees with the map this repo has used all along."""
    import yaml

    from edufair_renderer.template_inspect import inspect_template, to_layout_map

    generated = yaml.safe_load(to_layout_map(inspect_template(EXAMPLES / "template.pptx")))
    committed = yaml.safe_load((EXAMPLES / "layout-map.yaml").read_text(encoding="utf-8"))
    for key, entry in committed.items():
        assert key in generated, f"inspector lost layout {key}"
        assert generated[key]["layout_name"] == entry["layout_name"]
        assert generated[key]["regions"] == entry["regions"], key


def test_left_and_right_come_from_geometry_not_index_order(tmp_path):
    """A designer may ship Two Content with the right pane at the lower idx.

    Binding by index order would mirror every split slide silently.
    """
    from pptx import Presentation

    from edufair_renderer.template_inspect import _by_x, _r_two_content, Placeholder

    # idx 1 sits on the RIGHT, idx 2 on the left.
    title = Placeholder(0, None, 0, 0, 100, 10)
    right_at_idx1 = Placeholder(1, None, 5_000_000, 1_000_000, 4_000_000, 3_000_000)
    left_at_idx2 = Placeholder(2, None, 500_000, 1_000_000, 4_000_000, 3_000_000)
    bound = _r_two_content(title, [right_at_idx1, left_at_idx2])
    assert bound["left"] == 2 and bound["right"] == 1
    assert [p.idx for p in _by_x([right_at_idx1, left_at_idx2])] == [2, 1]


def test_missing_picture_placeholder_is_reported_not_silently_bound():
    from edufair_renderer.template_inspect import LayoutResult, report

    result = LayoutResult(
        key="Picture",
        layout_name="Picture with Caption",
        regions={"title": 0, "caption": 2},
        missing=["picture"],
        notes=["no PICTURE placeholder: images here would be free-positioned, not template-bound"],
    )
    text = report([result])
    assert "PARTIAL" in text and "picture" in text and "free-positioned" in text


# --- injecting the Cards layout into someone else's template -------------


def _theme_accents(pptx_path: Path) -> dict[str, str]:
    with zipfile.ZipFile(pptx_path) as z:
        theme = next(n for n in z.namelist() if n.startswith("ppt/theme/theme"))
        scheme = etree.fromstring(z.read(theme)).find(f".//{{{A_NS}}}clrScheme")
    return {
        f"accent{i}": scheme.find(f"{{{A_NS}}}accent{i}/{{{A_NS}}}srgbClr").get("val")
        for i in range(1, 5)
    }


def _card_tab_slots(pptx_path: Path) -> list[str]:
    """The theme slot each card tab fills from, left to right."""
    with zipfile.ZipFile(pptx_path) as z:
        for name in sorted(z.namelist()):
            if not name.startswith("ppt/slideLayouts/slideLayout"):
                continue
            root = etree.fromstring(z.read(name))
            if root.find(f".//{{{P_NS}}}cSld").get("name") != "Cards":
                continue
            slots = []
            for sp in root.findall(f".//{{{P_NS}}}sp"):
                label = sp.find(f".//{{{P_NS}}}cNvPr").get("name") or ""
                if "Tab" not in label:
                    continue
                clr = sp.find(f".//{{{P_NS}}}spPr/{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
                slots.append(clr.get("val"))
            return slots
    return []


def _rebrand(src: Path, dst: Path, accents: dict[str, str]) -> Path:
    """Copy a template with its theme accent colours replaced."""
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/theme/theme"):
                xml = data.decode("utf-8")
                for slot, value in accents.items():
                    pattern = r'(<a:' + slot + r'>)\s*<a:srgbClr val="[0-9A-Fa-f]{6}"/>'
                    xml = re.sub(pattern, r'\1<a:srgbClr val="' + value + '"/>', xml)
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    return dst


def test_cards_injects_into_a_stock_template(tmp_path):
    from edufair_renderer.cards_layout import inject_cards
    from edufair_renderer.template_inspect import inspect_template

    out = tmp_path / "with-cards.pptx"
    assert inject_cards(STOCK, out)["added"] is True

    cards = {r.key: r for r in inspect_template(out)}["Cards"]
    assert cards.layout_name == "Cards"
    assert not cards.missing
    assert [cards.regions[f"head{n}"] for n in range(1, 5)] == [1, 2, 3, 4]
    assert [cards.regions[f"card{n}"] for n in range(1, 5)] == [5, 6, 7, 8]


def test_injected_cards_take_the_target_theme_colours(tmp_path):
    """The point of injection: theme slots travel, colour values do not."""
    from edufair_renderer.cards_layout import inject_cards

    brand = {
        "accent1": "E4572E",
        "accent2": "17BEBB",
        "accent3": "FFC914",
        "accent4": "2E282A",
    }
    rebranded = _rebrand(STOCK, tmp_path / "rebranded.pptx", brand)
    inject_cards(STOCK, tmp_path / "plain-cards.pptx")
    inject_cards(rebranded, tmp_path / "brand-cards.pptx")

    slots = ["accent1", "accent2", "accent3", "accent4"]
    assert _card_tab_slots(tmp_path / "plain-cards.pptx") == slots
    assert _card_tab_slots(tmp_path / "brand-cards.pptx") == slots

    # Identical references, different palettes: the tabs resolve to each
    # template's own accents rather than anything we shipped.
    assert _theme_accents(tmp_path / "brand-cards.pptx") == brand
    assert _theme_accents(tmp_path / "plain-cards.pptx") != brand


def test_cards_injection_is_idempotent(tmp_path):
    from pptx import Presentation

    from edufair_renderer.cards_layout import inject_cards

    assert inject_cards(STOCK, tmp_path / "one.pptx")["added"] is True
    second = inject_cards(tmp_path / "one.pptx", tmp_path / "two.pptx")
    assert second["added"] is False and second["layout"] == "Cards"

    names = [
        layout.name
        for master in Presentation(str(tmp_path / "two.pptx")).slide_masters
        for layout in master.slide_layouts
    ]
    assert names.count("Cards") == 1


def test_cards_geometry_fits_a_short_widescreen_slide(tmp_path):
    """Fixed inch constants would push the panels off a 10 x 5.63 slide."""
    from pptx import Presentation
    from pptx.util import Inches

    from edufair_renderer.cards_layout import inject_cards

    wide = tmp_path / "wide.pptx"
    prs = Presentation(str(STOCK))
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.63)
    prs.save(str(wide))

    assert inject_cards(wide, tmp_path / "wide-cards.pptx")["added"] is True
    prs = Presentation(str(tmp_path / "wide-cards.pptx"))
    layout = next(
        lo for m in prs.slide_masters for lo in m.slide_layouts if lo.name == "Cards"
    )
    for ph in layout.placeholders:
        if ph.placeholder_format.idx in range(1, 9):
            assert ph.top + ph.height <= prs.slide_height, "card overflows the slide"


def test_injected_cards_render_a_real_session(tmp_path):
    from edufair_renderer.cards_layout import inject_cards
    from edufair_renderer.template_inspect import inspect_template, to_layout_map

    template = tmp_path / "with-cards.pptx"
    inject_cards(STOCK, template)
    map_path = tmp_path / "layout-map.yaml"
    map_path.write_text(to_layout_map(inspect_template(template)), encoding="utf-8")

    result = render_session(
        session_path=REPO / "examples" / "layout-gallery" / "gallery.md",
        template_path=template,
        layout_map_path=map_path,
        out_dir=tmp_path / "out",
    )

    cards_slide = None
    with zipfile.ZipFile(result["pptx"]) as z:
        for n in range(1, 11):
            xml = z.read(f"ppt/slides/slide{n}.xml").decode("utf-8")
            if "Tab is head1" in xml:
                cards_slide = etree.fromstring(xml.encode("utf-8"))
                break
    assert cards_slide is not None, "the Cards slide did not render"

    assert len(cards_slide.findall(f".//{{{P_NS}}}ph")) == 9
    # Every card is inherited: the slide contributes no styling of its own.
    for sp in cards_slide.findall(f".//{{{P_NS}}}sp"):
        if sp.find(f".//{{{P_NS}}}ph") is None:
            continue
        spPr = sp.find(f"{{{P_NS}}}spPr")
        assert spPr.find(f"{{{A_NS}}}solidFill") is None
        assert spPr.find(f"{{{A_NS}}}xfrm") is None


def test_regenerating_the_template_reproduces_the_committed_cards():
    """Geometry is exact integer fractions, so the template stays stable."""
    from pptx import Presentation

    from edufair_renderer.cards_layout import build_card_shapes

    committed = []
    with zipfile.ZipFile(EXAMPLES / "template.pptx") as z:
        for name in sorted(z.namelist()):
            if not name.startswith("ppt/slideLayouts/slideLayout"):
                continue
            root = etree.fromstring(z.read(name))
            if root.find(f".//{{{P_NS}}}cSld").get("name") != "Cards":
                continue
            committed = [
                sp
                for sp in root.findall(f".//{{{P_NS}}}sp")
                if "Card" in (sp.find(f".//{{{P_NS}}}cNvPr").get("name") or "")
            ]
            break
    assert committed, "no Cards layout in the committed template"

    prs = Presentation(str(EXAMPLES / "template.pptx"))
    rebuilt = build_card_shapes(prs.slide_width, prs.slide_height)

    # Compare geometry, not serialised bytes: a saved layout declares its
    # namespaces at the root, so its <p:sp> never matches a freshly built
    # element's serialisation even when the shapes are identical.
    def geometry(el):
        off = el.find(f".//{{{A_NS}}}off")
        ext = el.find(f".//{{{A_NS}}}ext")
        return (off.get("x"), off.get("y"), ext.get("cx"), ext.get("cy"))

    assert [geometry(el) for el in rebuilt] == [geometry(el) for el in committed]


# --- the course recipe and its blocks ------------------------------------


def _make_library(root: Path, *, resources: bool = True) -> Path:
    """A minimal course: one block, a deck, and non-slide resources."""
    import shutil

    root.mkdir(parents=True, exist_ok=True)
    shutil.copy(EXAMPLES / "template.pptx", root / "template.pptx")
    shutil.copy(EXAMPLES / "layout-map.yaml", root / "layout-map.yaml")

    block = root / "blocks" / "01-foundations"
    (block / "media").mkdir(parents=True)
    shutil.copy(
        EXAMPLES / "sessions" / "session-02.md", block / "slides.md"
    )

    meta = {
        "title": "Foundations of clinical teaching",
        "duration_minutes": 180,
        "competencies": {"CE1": "Designing clinical teaching sessions"},
        "resources": [],
    }
    if resources:
        (block / "lessonplan.md").write_text("# Lesson plan\n", encoding="utf-8")
        (block / "questions.xml").write_text("<questions/>\n", encoding="utf-8")
        meta["resources"] = [
            {"type": "lessonplan", "title": "Facilitator lesson plan", "path": "lessonplan.md"},
            {"type": "questions", "title": "Assessment bank", "path": "questions.xml"},
        ]
    (block / "block.yaml").write_text(
        yaml.safe_dump(meta, sort_keys=False, allow_unicode=True), encoding="utf-8"
    )

    (root / "course.yaml").write_text(
        yaml.safe_dump(
            {
                "id": "clinical-educator",
                "title": "Clinical Educator",
                "structure": [
                    {
                        "kind": "module",
                        "title": "Foundations",
                        "children": [
                            {
                                "kind": "day",
                                "title": "Day 1",
                                "children": [{"block": "01-foundations"}],
                            }
                        ],
                    }
                ],
            },
            sort_keys=False,
            allow_unicode=True,
        ),
        encoding="utf-8",
    )
    return root


def test_library_loads_course_blocks_and_resources(tmp_path):
    from edufair_renderer.library import load_library

    lib = load_library(_make_library(tmp_path / "lib"))
    assert lib.title == "Clinical Educator"
    assert list(lib.blocks) == ["01-foundations"]

    block = lib.blocks["01-foundations"]
    assert block.has_slides
    assert block.duration_minutes == 180
    assert [r.type for r in block.resources] == ["lessonplan", "questions"]

    # Arbitrary container depth under the course's own words.
    assert lib.structure[0]["kind"] == "module"
    assert lib.structure[0]["children"][0]["kind"] == "day"
    assert lib.structure[0]["children"][0]["children"][0] == {
        "kind": "block",
        "block": "01-foundations",
    }


def test_library_rejects_a_reference_to_a_missing_block(tmp_path):
    from edufair_renderer.library import LibraryError, load_library

    root = _make_library(tmp_path / "lib")
    course = yaml.safe_load((root / "course.yaml").read_text(encoding="utf-8"))
    course["structure"] = [{"block": "99-does-not-exist"}]
    (root / "course.yaml").write_text(yaml.safe_dump(course), encoding="utf-8")

    with pytest.raises(LibraryError, match="unknown block"):
        load_library(root)


def test_library_rejects_a_resource_that_is_not_there(tmp_path):
    from edufair_renderer.library import LibraryError, load_library

    root = _make_library(tmp_path / "lib")
    meta_path = root / "blocks" / "01-foundations" / "block.yaml"
    meta = yaml.safe_load(meta_path.read_text(encoding="utf-8"))
    meta["resources"].append({"type": "workbook", "path": "workbook.md"})
    meta_path.write_text(yaml.safe_dump(meta), encoding="utf-8")

    with pytest.raises(LibraryError, match="no such resource"):
        load_library(root)


def test_a_block_missing_from_the_recipe_is_unplaced_not_an_error(tmp_path):
    """Drafting a block before placing it must not break the build."""
    from edufair_renderer.library import load_library

    root = _make_library(tmp_path / "lib")
    extra = root / "blocks" / "02-draft"
    extra.mkdir()
    (extra / "slides.md").write_text(
        (EXAMPLES / "sessions" / "session-03.md").read_text(encoding="utf-8"),
        encoding="utf-8",
    )
    lib = load_library(root)
    assert "02-draft" in lib.blocks
    assert lib.structure[-1]["title"] == "Unplaced blocks"


def test_corpus_from_a_course_carries_structure_and_resources(tmp_path):
    from edufair_renderer.corpus import build_corpus

    root = _make_library(tmp_path / "lib")
    out = tmp_path / "out"
    build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    catalog = json.loads((out / "catalog.json").read_text())

    assert catalog["course"]["id"] == "clinical-educator"
    assert catalog["structure"][0]["kind"] == "module"

    block = catalog["blocks"][0]
    assert block["blockId"] == "01-foundations"
    assert block["durationMinutes"] == 180
    assert [r["type"] for r in block["resources"]] == ["lessonplan", "questions"]

    # Resources are copied beside the catalog, never parsed.
    for resource in block["resources"]:
        assert (out / resource["path"]).is_file()

    # Slides still carry their block, so the pane can group them.
    assert all(s["blockId"] == "01-foundations" for s in catalog["slides"])


def test_flat_sessions_still_build(tmp_path):
    """The legacy shape has to keep working while libraries migrate."""
    from edufair_renderer.corpus import build_corpus

    out = tmp_path / "out"
    summary = build_corpus(
        sessions_dir=EXAMPLES / "sessions",
        template=EXAMPLES / "template.pptx",
        layout_map=EXAMPLES / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    assert summary["sessions"] == 3
    catalog = json.loads((out / "catalog.json").read_text())
    assert catalog["structure"] == []
    assert catalog["blocks"] == []


# --- migration -----------------------------------------------------------


def _flat_library(root: Path) -> Path:
    import shutil

    (root / "sessions").mkdir(parents=True)
    for name in ("session-01.md", "session-02.md"):
        shutil.copy(EXAMPLES / "sessions" / name, root / "sessions" / name)
    shutil.copytree(EXAMPLES / "sessions" / "assets", root / "sessions" / "assets")
    shutil.copy(EXAMPLES / "template.pptx", root / "template.pptx")
    shutil.copy(EXAMPLES / "layout-map.yaml", root / "layout-map.yaml")

    (root / "credentials").mkdir()
    (root / "credentials" / "c.yaml").write_text(
        yaml.safe_dump(
            {
                "id": "cold-chain",
                "title": "Cold Chain",
                "modules": [
                    {"title": "Fundamentals", "sessions": ["01"]},
                    {"title": "Practice", "sessions": ["02"]},
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )
    return root


def test_migration_is_a_dry_run_until_asked(tmp_path):
    from edufair_renderer.migrate import plan_migration

    root = _flat_library(tmp_path / "lib")
    plan = plan_migration(root)
    assert len(plan["moves"]) == 2
    # Nothing written by planning alone.
    assert not (root / "course.yaml").exists()
    assert not (root / "blocks").exists()


def test_migration_builds_blocks_and_keeps_delivery_order(tmp_path):
    from edufair_renderer.corpus import build_corpus
    from edufair_renderer.library import load_library
    from edufair_renderer.migrate import apply_migration, plan_migration

    root = _flat_library(tmp_path / "lib")
    apply_migration(plan_migration(root))

    lib = load_library(root)
    # Identity comes from the credential, not the folder name.
    assert lib.course["id"] == "cold-chain"
    assert len(lib.blocks) == 2
    assert all(b.has_slides for b in lib.blocks.values())

    # The credential's modules survive as the course structure, at the top
    # level: the credential's own title names the course, not a container.
    assert [n["title"] for n in lib.structure] == ["Fundamentals", "Practice"]
    assert [n["kind"] for n in lib.structure] == ["module", "module"]
    assert lib.structure[0]["children"][0]["kind"] == "block"

    # Media referenced by a session moved into that block and still resolves.
    with_image = next(
        b for b in lib.blocks.values() if "media/" in b.slides.read_text(encoding="utf-8")
    )
    assert list((with_image.directory / "media").iterdir()), "media not copied"

    # And the migrated library renders.
    out = tmp_path / "out"
    summary = build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    assert summary["sessions"] == 2


def test_migration_refuses_a_library_already_migrated(tmp_path):
    from edufair_renderer.migrate import MigrationError, apply_migration, plan_migration

    root = _flat_library(tmp_path / "lib")
    apply_migration(plan_migration(root))
    with pytest.raises(MigrationError, match="already exists"):
        plan_migration(root)


# --- inline marks --------------------------------------------------------


def _marks(spec):
    return {
        name
        for name in ("bold", "italic", "superscript", "subscript", "strike", "underline", "code")
        if getattr(spec, name)
    }


def test_every_mark_parses():
    from edufair_renderer.runs import parse_emphasis

    runs = parse_emphasis("**b** *i* ***bi*** x^2^ H~2~O ~~s~~ __u__ `c`")
    by_text = {r.text: _marks(r) for r in runs if r.text.strip()}
    assert by_text["b"] == {"bold"}
    assert by_text["i"] == {"italic"}
    assert by_text["bi"] == {"bold", "italic"}
    assert by_text["s"] == {"strike"}
    assert by_text["u"] == {"underline"}
    assert by_text["c"] == {"code"}
    # Both "2" runs exist: one superscript, one subscript.
    twos = [_marks(r) for r in runs if r.text == "2"]
    assert {"superscript"} in twos and {"subscript"} in twos


def test_marks_nest():
    """A split regex cannot express this; the tokenizer must."""
    from edufair_renderer.runs import parse_emphasis, plain_text

    runs = parse_emphasis("**CO~2~ uptake**")
    assert all(r.bold for r in runs), "bold must survive across the nested mark"
    assert _marks(next(r for r in runs if r.text == "2")) == {"bold", "subscript"}
    assert plain_text("**CO~2~ uptake**") == "CO2 uptake"


def test_tilde_pairs_before_single_tilde():
    """Otherwise strikethrough would swallow subscript."""
    from edufair_renderer.runs import parse_emphasis

    runs = parse_emphasis("~~struck~~ and H~2~O")
    assert _marks(next(r for r in runs if r.text == "struck")) == {"strike"}
    assert _marks(next(r for r in runs if r.text == "2")) == {"subscript"}


def test_code_spans_are_literal():
    from edufair_renderer.runs import parse_emphasis

    (run,) = [r for r in parse_emphasis("`**not bold**`") if r.text.strip()]
    assert run.code and not run.bold
    assert run.text == "**not bold**"


def test_unclosed_and_underscores_stay_literal():
    """Authors write snake_case and stray asterisks; neither is a mark."""
    from edufair_renderer.runs import parse_emphasis, plain_text

    assert plain_text("unclosed ** stays") == "unclosed ** stays"
    assert plain_text("snake_case_name") == "snake_case_name"
    assert len(parse_emphasis("snake_case_name")) == 1


def test_marks_reach_the_slide_xml(tmp_path):
    """The attributes are the contract; the syntax is just how we get there."""
    import shutil

    lib = tmp_path / "lib"
    (lib / "sessions").mkdir(parents=True)
    shutil.copy(EXAMPLES / "template.pptx", lib / "template.pptx")
    shutil.copy(EXAMPLES / "layout-map.yaml", lib / "layout-map.yaml")
    (lib / "sessions" / "session-99.md").write_text(
        "---\n"
        'session: "99"\n'
        "title: Marks\n"
        "version: 1.0.0\n"
        "---\n\n"
        "--- slide\n"
        "id: m-01\n"
        "layout: Full\n"
        "title: Marks\n"
        "full:\n"
        "  type: p\n"
        '  text: "H~2~O and x^2^ and ~~gone~~ and __u__ and `kod`"\n'
        "---\n",
        encoding="utf-8",
    )
    result = render_session(
        session_path=lib / "sessions" / "session-99.md",
        template_path=lib / "template.pptx",
        layout_map_path=lib / "layout-map.yaml",
        out_dir=tmp_path / "out",
    )
    with zipfile.ZipFile(result["pptx"]) as z:
        xml = z.read("ppt/slides/slide1.xml").decode("utf-8")

    assert 'baseline="-25000"' in xml, "subscript"
    assert 'baseline="30000"' in xml, "superscript"
    assert 'strike="sngStrike"' in xml, "strikethrough"
    assert 'u="sng"' in xml, "underline"
    assert 'typeface="Consolas"' in xml, "inline code typeface"


def test_code_typeface_is_owned_by_the_library(tmp_path):
    """OOXML has no monospace theme slot, so the library names the font."""
    import shutil

    lib = tmp_path / "lib"
    (lib / "sessions").mkdir(parents=True)
    shutil.copy(EXAMPLES / "template.pptx", lib / "template.pptx")
    layout_map = (EXAMPLES / "layout-map.yaml").read_text(encoding="utf-8")
    (lib / "layout-map.yaml").write_text(
        "_style:\n  code_typeface: JetBrains Mono\n" + layout_map, encoding="utf-8"
    )
    (lib / "sessions" / "session-98.md").write_text(
        "---\n"
        'session: "98"\n'
        "title: Code\n"
        "version: 1.0.0\n"
        "---\n\n"
        "--- slide\n"
        "id: c-01\n"
        "layout: Full\n"
        "title: Code\n"
        "full:\n"
        "  type: p\n"
        '  text: "run `edufair-corpus` now"\n'
        "---\n",
        encoding="utf-8",
    )
    result = render_session(
        session_path=lib / "sessions" / "session-98.md",
        template_path=lib / "template.pptx",
        layout_map_path=lib / "layout-map.yaml",
        out_dir=tmp_path / "out",
    )
    with zipfile.ZipFile(result["pptx"]) as z:
        xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'typeface="JetBrains Mono"' in xml
    assert "Consolas" not in xml


def test_style_block_is_not_mistaken_for_a_layout(tmp_path):
    """_style sits in layout-map.yaml; the loader must skip it."""
    from edufair_renderer.layoutmap import load_layout_map, load_style

    path = tmp_path / "layout-map.yaml"
    path.write_text(
        "_style:\n  code_typeface: Fira Code\n"
        + (EXAMPLES / "layout-map.yaml").read_text(encoding="utf-8"),
        encoding="utf-8",
    )
    bindings = load_layout_map(path)
    assert "_style" not in bindings
    assert "Split" in bindings
    assert load_style(path)["code_typeface"] == "Fira Code"


def test_index_title_still_strips_every_marker(tmp_path):
    """plain_text feeds the semantic index; markers must never leak in."""
    import shutil

    lib = tmp_path / "lib"
    (lib / "sessions").mkdir(parents=True)
    shutil.copy(EXAMPLES / "template.pptx", lib / "template.pptx")
    shutil.copy(EXAMPLES / "layout-map.yaml", lib / "layout-map.yaml")
    (lib / "sessions" / "session-97.md").write_text(
        "---\n"
        'session: "97"\n'
        "title: T\n"
        "version: 1.0.0\n"
        "---\n\n"
        "--- slide\n"
        "id: t-01\n"
        "layout: Full\n"
        'title: "**CO~2~** in `air`"\n'
        "full:\n"
        "  type: p\n"
        "  text: body\n"
        "---\n",
        encoding="utf-8",
    )
    result = render_session(
        session_path=lib / "sessions" / "session-97.md",
        template_path=lib / "template.pptx",
        layout_map_path=lib / "layout-map.yaml",
        out_dir=tmp_path / "out",
    )
    entry = json.loads(result["index"].read_text())["slides"][0]
    assert entry["title"] == "CO2 in air"


# --- curriculum in SQLite, completion against it -------------------------


def _catalog_for_tracking() -> dict:
    return {
        "course": {"id": "clinical-educator", "title": "Clinical Educator",
                   "description": None},
        "structure": [],
        "blocks": [
            {"blockId": "01-foundations", "title": "Foundations",
             "durationMinutes": 180, "pptx": "sessions/01/x.pptx",
             "resources": [{"type": "workbook", "title": "Workbook",
                            "path": "blocks/01-foundations/workbook.md"}]},
            {"blockId": "02-feedback", "title": "Feedback",
             "durationMinutes": 180, "pptx": "sessions/02/x.pptx", "resources": []},
        ],
        "sessions": [],
        "competencies": {"CE1": "Designing", "CE2": "Teaching", "CE3": "Assessing"},
        "slides": [
            {"slideId": "a", "blockId": "01-foundations", "title": "A",
             "dok": 3, "develops": ["CE1"]},
            {"slideId": "b", "blockId": "01-foundations", "title": "B",
             "dok": 2, "develops": ["CE2"]},
            {"slideId": "c", "blockId": "02-feedback", "title": "C",
             "dok": 3, "develops": ["CE3"]},
        ],
        "credentials": [
            {"id": "ce", "title": "Clinical Educator", "ects": 5,
             "requires": [
                 {"competency": "CE1", "min_dok": 3},
                 {"competency": "CE2", "min_dok": 3},
                 {"competency": "CE3", "min_dok": 3},
             ]},
        ],
        "attribution": None,
    }


def test_sync_projects_the_curriculum(tmp_path):
    from edufair_renderer.tracking import connect, sync_curriculum

    conn = connect(tmp_path / "c.db")
    counts = sync_curriculum(conn, _catalog_for_tracking())
    assert counts["blocks"] == 2
    assert counts["slides"] == 3
    assert counts["credentials"] == 1
    assert counts["resources"] == 1
    conn.close()


def test_resync_is_idempotent_and_never_touches_events(tmp_path):
    """Curriculum is derived and rebuildable; events are the store of record."""
    from edufair_renderer.tracking import connect, record, sync_curriculum

    conn = connect(tmp_path / "c.db")
    catalog = _catalog_for_tracking()
    sync_curriculum(conn, catalog)
    record(conn, "ada", "completed", block_id="01-foundations")

    # Content changed upstream: a slide is retitled and a block renamed.
    catalog["blocks"][0]["title"] = "Foundations, revised"
    sync_curriculum(conn, catalog)

    assert conn.execute("SELECT COUNT(*) FROM blocks").fetchone()[0] == 2, "duplicated"
    assert conn.execute(
        "SELECT title FROM blocks WHERE block_id='01-foundations'"
    ).fetchone()[0] == "Foundations, revised"
    assert conn.execute("SELECT COUNT(*) FROM events").fetchone()[0] == 1, "events lost"
    conn.close()


def test_completion_evidences_competencies_at_the_slides_dok(tmp_path):
    from edufair_renderer.tracking import connect, record, sync_curriculum

    conn = connect(tmp_path / "c.db")
    sync_curriculum(conn, _catalog_for_tracking())
    record(conn, "ada", "completed", block_id="01-foundations")

    evidence = {
        row["competency_id"]: row["dok"]
        for row in conn.execute(
            "SELECT competency_id, dok FROM learner_competency WHERE learner='ada'"
        )
    }
    assert evidence == {"CE1": 3, "CE2": 2}
    conn.close()


def test_credential_requires_the_dok_not_just_the_competency(tmp_path):
    """DOK 2 evidence must not satisfy a DOK 3 requirement."""
    from edufair_renderer.tracking import (
        connect,
        credential_progress,
        record,
        sync_curriculum,
    )

    conn = connect(tmp_path / "c.db")
    sync_curriculum(conn, _catalog_for_tracking())
    record(conn, "ada", "completed", block_id="01-foundations")
    record(conn, "ada", "completed", block_id="02-feedback")

    (entry,) = credential_progress(conn, "ada")
    met = {r["competency"]: r["met"] for r in entry["requirements"]}
    assert met == {"CE1": True, "CE2": False, "CE3": True}
    assert entry["complete"] is False, "CE2 is only evidenced at DOK 2"
    conn.close()


def test_a_learner_with_no_events_has_no_evidence(tmp_path):
    from edufair_renderer.tracking import connect, credential_progress, sync_curriculum

    conn = connect(tmp_path / "c.db")
    sync_curriculum(conn, _catalog_for_tracking())
    (entry,) = credential_progress(conn, "nobody")
    assert not any(r["met"] for r in entry["requirements"])
    assert entry["complete"] is False
    conn.close()


def test_completion_against_an_unknown_block_is_refused(tmp_path):
    """Otherwise the record accumulates events pointing at nothing."""
    from edufair_renderer.tracking import TrackingError, connect, record, sync_curriculum

    conn = connect(tmp_path / "c.db")
    sync_curriculum(conn, _catalog_for_tracking())
    with pytest.raises(TrackingError, match="no such block"):
        record(conn, "ada", "completed", block_id="99-does-not-exist")
    conn.close()


def test_sync_refuses_a_catalog_with_no_course(tmp_path):
    """A flat library predates the recipe and has nothing to hang events on."""
    from edufair_renderer.tracking import TrackingError, connect, sync_curriculum

    conn = connect(tmp_path / "c.db")
    with pytest.raises(TrackingError, match="course id"):
        sync_curriculum(conn, {"course": {}, "blocks": [], "slides": []})
    conn.close()


def test_tracking_syncs_a_really_rendered_catalog(tmp_path):
    """End to end: render a library, project it, record against it."""
    import shutil

    from edufair_renderer.corpus import build_corpus
    from edufair_renderer.tracking import connect, credential_progress, record, sync_curriculum

    root = _make_library(tmp_path / "lib")
    out = tmp_path / "out"
    build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    catalog = json.loads((out / "catalog.json").read_text())

    conn = connect(tmp_path / "c.db")
    counts = sync_curriculum(conn, catalog)
    assert counts["blocks"] == 1
    assert counts["resources"] == 2

    record(conn, "ada", "completed", block_id="01-foundations", source="pr")
    assert conn.execute(
        "SELECT COUNT(*) FROM learner_competency WHERE learner='ada'"
    ).fetchone()[0] >= 1
    # No credentials in this fixture, so progress is simply empty.
    assert credential_progress(conn, "ada") == []
    conn.close()


def test_check_blocks_is_quiet_on_a_good_library(tmp_path):
    from edufair_renderer.library import check_blocks

    assert check_blocks(_make_library(tmp_path / "lib")) == []


def test_check_blocks_wants_a_lesson_plan(tmp_path):
    """A deck with no lesson plan is a slideshow, not a session."""
    from edufair_renderer.library import check_blocks

    root = _make_library(tmp_path / "lib", resources=False)
    problems = check_blocks(root)
    assert [p["where"] for p in problems] == ["01-foundations"]
    assert "lessonplan.md" in problems[0]["message"]


def test_check_blocks_reports_rather_than_raises(tmp_path):
    """An author wants every problem at once, not just the first."""
    from edufair_renderer.library import check_blocks

    root = _make_library(tmp_path / "lib")
    (root / "blocks" / "01-foundations" / "lessonplan.md").unlink()
    meta_path = root / "blocks" / "01-foundations" / "block.yaml"
    meta = yaml.safe_load(meta_path.read_text(encoding="utf-8"))
    meta["resources"].append({"type": "workbook", "path": "nowhere.md"})
    meta_path.write_text(yaml.safe_dump(meta), encoding="utf-8")

    messages = [p["message"] for p in check_blocks(root)]
    assert any("lessonplan.md" in m for m in messages)
    assert any("nowhere.md" in m for m in messages)


def test_a_block_may_have_a_lesson_plan_and_no_deck(tmp_path):
    """Not every session is taught from slides."""
    from edufair_renderer.library import check_blocks, load_library

    root = _make_library(tmp_path / "lib")
    (root / "blocks" / "01-foundations" / "slides.md").unlink()

    block = load_library(root).blocks["01-foundations"]
    assert not block.has_slides
    assert [r.type for r in block.resources] == ["lessonplan", "questions"]
    assert check_blocks(root) == []


def _with_outcomes(root: Path, catalogue: dict, block_outcomes=None) -> Path:
    """Add an outcome catalogue to a library made by _make_library."""
    (root / "competencies").mkdir(exist_ok=True)
    (root / "competencies" / "framework.yaml").write_text(
        yaml.safe_dump(
            {"competencies": {"CE1": {"label": "Designing clinical teaching sessions"}}}
        ),
        encoding="utf-8",
    )
    (root / "outcomes.yaml").write_text(
        yaml.safe_dump({"outcomes": catalogue}, sort_keys=False), encoding="utf-8"
    )
    if block_outcomes is not None:
        meta_path = root / "blocks" / "01-foundations" / "block.yaml"
        meta = yaml.safe_load(meta_path.read_text(encoding="utf-8"))
        meta["outcomes"] = block_outcomes
        meta_path.write_text(yaml.safe_dump(meta, sort_keys=False), encoding="utf-8")
    return root


def test_an_outcome_may_not_develop_a_competency_that_does_not_exist(tmp_path):
    """A claim the library cannot support is an error, not a warning."""
    from edufair_renderer.outcomes import OutcomeError, load_outcomes

    root = _make_library(tmp_path / "lib")
    _with_outcomes(root, {"O1": {"statement": "Do a thing", "develops": ["NOPE"]}})

    with pytest.raises(OutcomeError, match="unknown competency"):
        load_outcomes(root / "outcomes.yaml", {"CE1"})


def test_a_slide_develops_what_its_outcomes_develop(tmp_path):
    """The derivation everything downstream depends on."""
    from edufair_renderer.outcomes import derive_develops, load_outcomes

    root = _make_library(tmp_path / "lib")
    _with_outcomes(root, {"O1": {"statement": "Teach well", "develops": ["CE1"]}})
    catalogue = load_outcomes(root / "outcomes.yaml", {"CE1"})

    assert derive_develops([], ["O1"], catalogue) == ["CE1"]
    # Declared competencies keep their place, so a library that has not
    # adopted outcomes renders byte-identically.
    assert derive_develops(["CE2"], ["O1"], catalogue) == ["CE2", "CE1"]
    assert derive_develops(["CE1"], ["O1"], catalogue) == ["CE1"]
    assert derive_develops(["CE1"], [], {}) == ["CE1"]


def test_a_bare_string_outcome_is_just_its_statement(tmp_path):
    from edufair_renderer.outcomes import load_outcomes

    root = _make_library(tmp_path / "lib")
    _with_outcomes(root, {"O1": "Not yet mapped to anything"})
    catalogue = load_outcomes(root / "outcomes.yaml", {"CE1"})
    assert catalogue["O1"].statement == "Not yet mapped to anything"
    assert catalogue["O1"].develops == []


def test_alignment_reports_every_kind_of_gap(tmp_path):
    from edufair_renderer.library import check_alignment

    root = _make_library(tmp_path / "lib")
    _with_outcomes(
        root,
        {
            "O1": {"statement": "Served by a slide", "develops": ["CE1"]},
            "O2": {"statement": "Declared but unserved", "develops": ["CE1"]},
            "O3": {"statement": "Addressed by nobody", "develops": ["CE1"]},
        },
        block_outcomes=["O1", "O2"],
    )
    slides = root / "blocks" / "01-foundations" / "slides.md"
    text = slides.read_text(encoding="utf-8")
    # First slide serves O1; second references an outcome that does not exist.
    text = text.replace("--- slide", "--- slide\noutcomes: [O1]", 1)
    text = text.replace("--- slide\nid", "--- slide\noutcomes: [O99]\nid", 1)
    slides.write_text(text, encoding="utf-8")

    problems = check_alignment(root)
    messages = [p["message"] for p in problems]
    errors = [p for p in problems if p["level"] == "error"]

    assert any("unknown outcome 'O99'" in m for m in messages)
    assert errors, "an unknown outcome reference is an error"
    assert any("'O2' but no slide serves it" in m for m in messages)
    assert any("'O3' belongs to no session" in m for m in messages)


def test_alignment_is_quiet_without_a_catalogue(tmp_path):
    """Additive: a library that never adopts outcomes is not nagged."""
    from edufair_renderer.library import check_alignment

    assert check_alignment(_make_library(tmp_path / "lib")) == []


def test_outcomes_on_a_slide_is_metadata_not_a_region(tmp_path):
    from edufair_renderer.parser import SessionParseError, parse_session

    path = tmp_path / "s.md"
    path.write_text(
        "---\nsession: '01'\ntitle: T\nversion: 1.0.0\n---\n\n"
        "--- slide\nid: s1\nlayout: Full\noutcomes: {text: not a list}\n---\n",
        encoding="utf-8",
    )
    with pytest.raises(SessionParseError, match="not a region"):
        parse_session(path)


def test_a_slide_with_only_outcomes_reaches_the_competency_everywhere(tmp_path):
    """The whole point, proved end to end.

    A slide names an outcome and no competency. That competency must
    still appear in catalog.json, in the slides projection, and in
    learner_competency — which is how we know nothing downstream had to
    learn about outcomes.
    """
    import sqlite3

    from edufair_renderer.corpus import build_corpus
    from edufair_renderer.tracking import connect, record, sync_curriculum

    root = _make_library(tmp_path / "lib")
    _with_outcomes(
        root,
        {"O1": {"statement": "Teach a clinical session well", "develops": ["CE1"], "dok": 3}},
        block_outcomes=["O1"],
    )
    slides = root / "blocks" / "01-foundations" / "slides.md"
    text = slides.read_text(encoding="utf-8")
    # Strip any competency the fixture declares, so CE1 can only arrive
    # through the outcome.
    text = re.sub(r"(?m)^develops:.*(?:\n[ \t]+-.*)*\n", "", text)
    text = text.replace("--- slide", "--- slide\noutcomes: [O1]", 1)
    slides.write_text(text, encoding="utf-8")

    out = tmp_path / "out"
    build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    catalog = json.loads((out / "catalog.json").read_text(encoding="utf-8"))

    assert catalog["outcomes"]["O1"]["develops"] == ["CE1"]
    first = catalog["slides"][0]
    assert first["outcomes"] == ["O1"]
    assert "CE1" in first["develops"], "the derived competency must reach catalog.json"

    conn = connect(tmp_path / "t.db")
    sync_curriculum(conn, catalog)
    assert conn.execute(
        "SELECT competency_id FROM outcome_competencies WHERE outcome_id='O1'"
    ).fetchone()[0] == "CE1"
    assert conn.execute("SELECT COUNT(*) FROM slide_outcomes").fetchone()[0] >= 1

    record(conn, "learner@example.org", "completed", block_id="01-foundations")
    competencies = {
        row[0] for row in conn.execute("SELECT competency_id FROM learner_competency")
    }
    assert "CE1" in competencies, "evidence must flow through the outcome"
    outcomes = {row[0] for row in conn.execute("SELECT outcome_id FROM learner_outcome")}
    assert outcomes == {"O1"}


def test_a_library_without_outcomes_renders_identically(tmp_path):
    """Additive: adopting outcomes is a choice, not a migration."""
    from edufair_renderer.corpus import build_corpus

    root = _make_library(tmp_path / "lib")
    out = tmp_path / "out"
    build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    catalog = json.loads((out / "catalog.json").read_text(encoding="utf-8"))
    assert catalog["outcomes"] == {}
    assert catalog["slides"][0]["outcomes"] == []


def test_lifting_outcomes_is_a_dry_run_until_asked(tmp_path):
    from edufair_renderer.migrate import plan_outcomes

    root = _make_library(tmp_path / "lib")
    plan = plan_outcomes(root)

    assert not (root / "outcomes.yaml").exists(), "planning must write nothing"
    # The fixture deck carries its outcomes as free text in frontmatter,
    # which is exactly the shape this lifts.
    assert [o["statement"] for o in plan["catalogue"].values()] == [
        "Maintain compliant cold chain documentation",
    ]
    assert plan["blocks"][0]["outcomes"] == ["O1"]
    # Seeded from the block's own competencies: a starting point the
    # printed plan is explicit about, not a finding.
    assert plan["catalogue"]["O1"]["develops"] == ["CE1"]


def test_lifting_outcomes_writes_the_catalogue_and_points_blocks_at_it(tmp_path):
    from edufair_renderer.migrate import apply_outcomes, plan_outcomes

    root = _make_library(tmp_path / "lib")
    apply_outcomes(plan_outcomes(root))

    catalogue = yaml.safe_load((root / "outcomes.yaml").read_text(encoding="utf-8"))
    assert (
        catalogue["outcomes"]["O1"]["statement"]
        == "Maintain compliant cold chain documentation"
    )
    meta = yaml.safe_load(
        (root / "blocks" / "01-foundations" / "block.yaml").read_text(encoding="utf-8")
    )
    assert meta["outcomes"] == ["O1"]
    # Everything the block already declared survives.
    assert meta["duration_minutes"] == 180


def test_the_same_wording_in_two_sessions_makes_two_outcomes(tmp_path):
    """An outcome belongs to one session, so wording is not merged.

    Merging would quietly turn a session outcome into something that
    spans sessions, which is a competency and a different thing.
    """
    import shutil

    from edufair_renderer.migrate import plan_outcomes

    root = _make_library(tmp_path / "lib")
    shutil.copytree(root / "blocks" / "01-foundations", root / "blocks" / "02-more")

    plan = plan_outcomes(root)
    assert len(plan["catalogue"]) == 2, "two sessions, two outcomes"
    assert [b["outcomes"] for b in plan["blocks"]] == [["O1"], ["O2"]]
    # And the repeat is reported, because it usually means the two
    # sessions need distinguishing.
    assert plan["repeated"], "repeated wording must be surfaced"
    statement, first, again = plan["repeated"][0]
    assert first == "01-foundations" and again == "02-more"


def test_lifting_outcomes_refuses_to_overwrite_a_catalogue(tmp_path):
    from edufair_renderer.migrate import MigrationError, plan_outcomes

    root = _make_library(tmp_path / "lib")
    (root / "outcomes.yaml").write_text("outcomes: {}\n", encoding="utf-8")
    with pytest.raises(MigrationError, match="already exists"):
        plan_outcomes(root)


def test_an_outcome_belongs_to_one_session(tmp_path):
    """Two sessions claiming one outcome is a modelling mistake.

    An outcome is what a session is for. Something that spans sessions is
    a competency, and calling it an outcome hides that.
    """
    import shutil

    from edufair_renderer.library import check_alignment

    root = _make_library(tmp_path / "lib")
    shutil.copytree(root / "blocks" / "01-foundations", root / "blocks" / "02-more")
    _with_outcomes(
        root,
        {"O1": {"statement": "Shared by two", "develops": ["CE1"]}},
        block_outcomes=["O1"],
    )
    second = root / "blocks" / "02-more" / "block.yaml"
    meta = yaml.safe_load(second.read_text(encoding="utf-8"))
    meta["outcomes"] = ["O1"]
    second.write_text(yaml.safe_dump(meta), encoding="utf-8")

    problems = check_alignment(root)
    clash = [p for p in problems if "belongs to" in p["message"]]
    assert clash and clash[0]["level"] == "error"
    assert "something that spans sessions is a competency" in clash[0]["message"].lower()


def test_a_competency_confined_to_one_session_is_flagged(tmp_path):
    """A competency builds across sessions; one that does not is an outcome."""
    from edufair_renderer.library import check_alignment

    root = _make_library(tmp_path / "lib")
    _with_outcomes(
        root,
        {"O1": {"statement": "Only ever here", "develops": ["CE1"]}},
        block_outcomes=["O1"],
    )
    messages = [p["message"] for p in check_alignment(root)]
    assert any("built only in" in m for m in messages)
    # A gap, not a failure: plenty of courses have one session per topic.
    assert all(
        p["level"] == "warning"
        for p in check_alignment(root)
        if "built only in" in p["message"]
    )


def test_a_competency_across_two_sessions_is_not_flagged(tmp_path):
    import shutil

    from edufair_renderer.library import check_alignment

    root = _make_library(tmp_path / "lib")
    shutil.copytree(root / "blocks" / "01-foundations", root / "blocks" / "02-more")
    _with_outcomes(
        root,
        {
            "O1": {"statement": "First step", "develops": ["CE1"]},
            "O2": {"statement": "Building on it", "develops": ["CE1"]},
        },
        block_outcomes=["O1"],
    )
    second = root / "blocks" / "02-more" / "block.yaml"
    meta = yaml.safe_load(second.read_text(encoding="utf-8"))
    meta["outcomes"] = ["O2"]
    second.write_text(yaml.safe_dump(meta), encoding="utf-8")

    messages = [p["message"] for p in check_alignment(root)]
    assert not any("built only in" in m for m in messages)


def test_a_competency_progression_shows_where_it_is_built(tmp_path):
    """The distinction, made visible rather than only checkable.

    An outcome belongs to one session. A competency runs through several,
    and the progression is the list of those sessions in delivery order.
    """
    import shutil

    from edufair_renderer.corpus import build_corpus
    from edufair_renderer.tracking import connect, sync_curriculum

    root = _make_library(tmp_path / "lib")
    shutil.copytree(root / "blocks" / "01-foundations", root / "blocks" / "02-more")
    _with_outcomes(
        root,
        {
            "O1": {"statement": "Recognise it", "develops": ["CE1"], "dok": 1},
            "O2": {"statement": "Do it unaided", "develops": ["CE1"], "dok": 3},
        },
        block_outcomes=["O1"],
    )
    second = root / "blocks" / "02-more" / "block.yaml"
    meta = yaml.safe_load(second.read_text(encoding="utf-8"))
    meta["outcomes"] = ["O2"]
    second.write_text(yaml.safe_dump(meta), encoding="utf-8")

    course_path = root / "course.yaml"
    course = yaml.safe_load(course_path.read_text(encoding="utf-8"))
    course["structure"] = [{"block": "01-foundations"}, {"block": "02-more"}]
    course_path.write_text(yaml.safe_dump(course), encoding="utf-8")

    out = tmp_path / "out"
    build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    catalog = json.loads((out / "catalog.json").read_text(encoding="utf-8"))

    # Each outcome records the session it belongs to.
    assert catalog["outcomes"]["O1"]["block"] == "01-foundations"
    assert catalog["outcomes"]["O2"]["block"] == "02-more"

    conn = connect(tmp_path / "t.db")
    sync_curriculum(conn, catalog)
    rows = conn.execute(
        "SELECT block_id, dok FROM competency_progression WHERE competency_id='CE1'"
    ).fetchall()
    assert [r[0] for r in rows] == ["01-foundations", "02-more"], "in delivery order"
    assert [r[1] for r in rows] == [1, 3], "and it builds"


def _role_deck(root: Path) -> Path:
    """Replace a block's deck with two role slides and nothing typed."""
    slides = root / "blocks" / "01-foundations" / "slides.md"
    slides.write_text(
        "---\n"
        "session: '01'\n"
        "title: Foundations\n"
        "version: 1.0.0\n"
        "---\n\n"
        "--- slide\n"
        "id: s-01\n"
        "layout: Title\n"
        "role: title\n"
        "---\n\n"
        "--- slide\n"
        "id: s-02\n"
        "layout: Full\n"
        "role: outcomes\n"
        "---\n",
        encoding="utf-8",
    )
    return slides


def test_a_role_slide_is_filled_from_the_library(tmp_path):
    """Nothing is typed twice: the deck says what a slide is, not what it says."""
    from edufair_renderer.corpus import build_corpus

    root = _make_library(tmp_path / "lib")
    _with_outcomes(
        root,
        {
            "O1": {"statement": "Plan a clinical teaching session", "develops": ["CE1"]},
            "O2": {"statement": "Debrief a learner afterwards", "develops": ["CE1"]},
        },
        block_outcomes=["O1", "O2"],
    )
    _role_deck(root)
    meta_path = root / "blocks" / "01-foundations" / "block.yaml"
    meta = yaml.safe_load(meta_path.read_text(encoding="utf-8"))
    meta["code"] = "CE01"
    meta_path.write_text(yaml.safe_dump(meta, sort_keys=False), encoding="utf-8")

    out = tmp_path / "out"
    build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    catalog = json.loads((out / "catalog.json").read_text(encoding="utf-8"))
    by_id = {s["slideId"]: s for s in catalog["slides"]}

    # The title slide carries the session's title, from block.yaml.
    assert by_id["s-01"]["title"] == "Foundations of clinical teaching"

    # The outcomes slide carries the statements, from outcomes.yaml, and
    # serves those outcomes without anyone tagging it.
    assert by_id["s-02"]["title"] == "Learning outcomes"
    assert by_id["s-02"]["outcomes"] == ["O1", "O2"]
    assert by_id["s-02"]["develops"] == ["CE1"], "derived through the outcomes"

    # And the words reached the rendered deck, not just the index.
    deck = out / "sessions" / "01-foundations"
    pptx = next(deck.glob("*.pptx"))
    with zipfile.ZipFile(pptx) as z:
        xml = z.read("ppt/slides/slide2.xml").decode("utf-8")
    assert "Plan a clinical teaching session" in xml
    assert "Debrief a learner afterwards" in xml


def test_what_an_author_writes_on_a_role_slide_wins(tmp_path):
    """A role fills what is absent; it never overwrites."""
    from edufair_renderer.corpus import build_corpus

    root = _make_library(tmp_path / "lib")
    _with_outcomes(
        root, {"O1": {"statement": "The catalogue wording", "develops": ["CE1"]}},
        block_outcomes=["O1"],
    )
    slides = _role_deck(root)
    slides.write_text(
        slides.read_text(encoding="utf-8").replace(
            "role: outcomes", "role: outcomes\ntitle: What you will be able to do"
        ),
        encoding="utf-8",
    )
    out = tmp_path / "out"
    build_corpus(
        library=root,
        template=root / "template.pptx",
        layout_map=root / "layout-map.yaml",
        out_dir=out,
        warn=lambda m: None,
    )
    catalog = json.loads((out / "catalog.json").read_text(encoding="utf-8"))
    by_id = {s["slideId"]: s for s in catalog["slides"]}
    assert by_id["s-02"]["title"] == "What you will be able to do"


def test_an_unknown_role_is_refused(tmp_path):
    from edufair_renderer.parser import SessionParseError, parse_session

    path = tmp_path / "s.md"
    path.write_text(
        "---\nsession: '01'\ntitle: T\nversion: 1.0.0\n---\n\n"
        "--- slide\nid: s1\nlayout: Full\nrole: whatever\n---\n",
        encoding="utf-8",
    )
    with pytest.raises(SessionParseError, match="unknown role"):
        parse_session(path)


# ---- how a picture meets its frame --------------------------------------


def _picture_deck(tmp_path, fit, size=(600, 1200)):
    """Render one slide with a picture in a Picture placeholder."""
    from PIL import Image

    from pptx import Presentation

    from edufair_renderer.render import render_session

    sessions = tmp_path / "sessions"
    (sessions / "assets").mkdir(parents=True)
    Image.new("RGB", size, (200, 60, 60)).save(sessions / "assets" / "p.png")
    fit_line = f"\n  fit: {fit}" if fit else ""
    (sessions / "s.md").write_text(
        "---\nsession: '01'\ntitle: T\nversion: 1.0.0\n---\n\n"
        "--- slide\nid: s1\nlayout: Picture\ntitle: A picture\n"
        f"picture:\n  type: image\n  src: assets/p.png{fit_line}\n---\n",
        encoding="utf-8",
    )
    out = tmp_path / "out"
    render_session(sessions / "s.md", EXAMPLES / "template.pptx",
                   EXAMPLES / "layout-map.yaml", out)

    prs = Presentation(str(next(out.glob("*.pptx"))))
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    pic = next(s for s in prs.slides[0].shapes if s.element.find(f".//{A}blip") is not None)
    rect = pic.element.find(f".//{A}srcRect")
    crop_top = int(rect.get("t") or 0) / 100000 if rect is not None else 0.0
    crop_left = int(rect.get("l") or 0) / 100000 if rect is not None else 0.0
    return {
        "width": pic.width,
        "height": pic.height,
        "crop_top": crop_top,
        "crop_left": crop_left,
        "bound": pic.element.find(f".//{P}ph") is not None,
    }


def test_a_picture_fills_or_fits_as_the_slide_asks(tmp_path):
    """A portrait picture in a landscape frame forces a choice.

    The template knows the shape of the hole; only the author knows
    whether this picture may lose its edges.
    """
    cover = _picture_deck(tmp_path / "a", "cover")
    contain = _picture_deck(tmp_path / "b", "contain")

    # Cover fills the frame and takes the overflow off both ends.
    assert cover["crop_top"] > 0.3
    assert cover["width"] > contain["width"]
    # Contain keeps every pixel and sits narrower than the frame.
    assert contain["crop_top"] == 0
    assert contain["height"] == cover["height"]

    for result in (cover, contain):
        assert result["bound"], "the slide must stay bound to its placeholder"


def test_matching_one_edge_is_consistent_whatever_the_picture(tmp_path):
    """`width` and `height` say which edge governs, rather than letting
    the picture's own shape decide as cover does."""
    tall = _picture_deck(tmp_path / "tall", "width", size=(600, 1200))
    wide = _picture_deck(tmp_path / "wide", "width", size=(2000, 500))

    # Both match the frame's width; only the tall one loses anything.
    assert tall["width"] == wide["width"]
    assert tall["crop_top"] > 0
    assert wide["crop_top"] == 0 and wide["crop_left"] == 0


def test_a_picture_with_no_fit_still_fills_the_frame(tmp_path):
    """The default is unchanged, so existing libraries render as before."""
    assert _picture_deck(tmp_path, None)["crop_top"] > 0.3


def test_an_unknown_fit_is_refused_at_parse_time(tmp_path):
    from edufair_renderer.parser import SessionParseError, parse_session

    path = tmp_path / "s.md"
    path.write_text(
        "---\nsession: '01'\ntitle: T\nversion: 1.0.0\n---\n\n"
        "--- slide\nid: s1\nlayout: Picture\n"
        "picture:\n  type: image\n  src: a.png\n  fit: squish\n---\n",
        encoding="utf-8",
    )
    with pytest.raises(SessionParseError, match="fit must be one of"):
        parse_session(path)


# ---- lists have to look like lists ------------------------------------
#
# A `ul` rendered without bullets is simply wrong: the content says the
# list is bulleted and the output does not show it. The example template
# has no `bodyStyle` element at all, so nothing defined a bullet and
# every bulleted list came out as plain paragraphs -- which is what was
# reported, and what nothing here would have caught.

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _bullet_props(paragraph):
    """The bullet elements written onto a paragraph, by name."""
    pPr = paragraph._p.find(f"{{{A_NS}}}pPr")
    if pPr is None:
        return []
    out = []
    for child in pPr:
        name = etree.QName(child).localname
        if name.startswith("bu"):
            out.append(name)
    return out


def _body_paragraphs(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
            continue  # the title
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text.strip():
                yield paragraph


def _render_one(tmp_path, region_yaml):
    lib = tmp_path / "lib"
    (lib / "blocks" / "01-b").mkdir(parents=True)
    shutil.copy(EXAMPLES / "template.pptx", lib / "template.pptx")
    shutil.copy(EXAMPLES / "layout-map.yaml", lib / "layout-map.yaml")
    (lib / "course.yaml").write_text(
        "title: T\nstructure:\n  - block: 01-b\n", encoding="utf-8"
    )
    (lib / "blocks" / "01-b" / "block.yaml").write_text(
        "title: B\nduration_minutes: 45\nresources: []\n", encoding="utf-8"
    )
    (lib / "blocks" / "01-b" / "slides.md").write_text(
        "---\nsession: '01'\ntitle: T\nversion: 1.0.0\n---\n\n"
        "--- slide\nid: s-01\nlayout: Full\ntitle: T\n" + region_yaml + "---\n",
        encoding="utf-8",
    )
    out = tmp_path / "build"
    # Through the CLI, which is what derives the template and layout
    # map from the library rather than being handed them.
    corpus_main([str(lib), "--out", str(out)])
    deck = next((out / "sessions").rglob("*.pptx"))
    return Presentation(str(deck)).slides[0]


def test_a_bulleted_list_gets_bullets_even_from_a_bare_template(tmp_path):
    slide = _render_one(
        tmp_path,
        "full:\n  type: ul\n  items:\n    - First\n"
        "    - text: Parent\n      items:\n        - Nested\n",
    )
    props = [_bullet_props(p) for p in _body_paragraphs(slide)]
    assert props, "no body paragraphs were written at all"
    for one in props:
        assert "buChar" in one, f"a bulleted item with no bullet: {one}"


def test_a_numbered_list_is_numbered_not_bulleted(tmp_path):
    slide = _render_one(tmp_path, "full:\n  type: ol\n  items:\n    - One\n    - Two\n")
    for one in (_bullet_props(p) for p in _body_paragraphs(slide)):
        assert "buAutoNum" in one
        assert "buChar" not in one


def test_a_placeholder_can_hold_both_plain_lines_and_bullets(tmp_path):
    """One region, a lead-in, the points, and a line that closes them.

    A layout offers the regions it offers; splitting a lead-in into its
    own region is not possible when there is only one, and inventing a
    second would be a layout decision taken by the content.
    """
    slide = _render_one(
        tmp_path,
        "full:\n  type: ul\n  items:\n"
        "    - text: Three things to check.\n      bullet: false\n"
        "    - Who can be re-identified\n"
        "    - What was consented to\n"
        "    - text: If any is unclear, do not share.\n      bullet: false\n",
    )
    got = [(p.text, _bullet_props(p)) for p in _body_paragraphs(slide)]
    assert [text for text, _ in got] == [
        "Three things to check.",
        "Who can be re-identified",
        "What was consented to",
        "If any is unclear, do not share.",
    ]
    assert "buNone" in got[0][1]
    assert "buChar" in got[1][1]
    assert "buChar" in got[2][1]
    assert "buNone" in got[3][1]


def test_a_template_that_defines_its_own_bullet_is_left_alone(tmp_path):
    """The glyph is the template's. Only its absence is ours to fix."""
    lib = tmp_path / "lib"
    (lib / "blocks" / "01-b").mkdir(parents=True)
    shutil.copy(EXAMPLES / "template.pptx", lib / "template.pptx")

    # Give the master a bodyStyle that asks for a dash at the first level.
    import zipfile

    src = lib / "template.pptx"
    data = {}
    with zipfile.ZipFile(src) as z:
        for name in z.namelist():
            data[name] = z.read(name)
    master = next(n for n in data if n.startswith("ppt/slideMasters/slideMaster1.xml"))
    root = etree.fromstring(data[master])
    styles = root.find(f"{{{A_NS}}}txStyles")
    if styles is None:
        styles = etree.SubElement(root, f"{{{A_NS}}}txStyles")
    body = etree.SubElement(styles, f"{{{A_NS}}}bodyStyle")
    lvl = etree.SubElement(body, f"{{{A_NS}}}lvl1pPr")
    etree.SubElement(lvl, f"{{{A_NS}}}buChar").set("char", "-")
    data[master] = etree.tostring(root, xml_declaration=True, encoding="UTF-8")
    with zipfile.ZipFile(src, "w") as z:
        for name, blob in data.items():
            z.writestr(name, blob)

    shutil.copy(EXAMPLES / "layout-map.yaml", lib / "layout-map.yaml")
    (lib / "course.yaml").write_text("title: T\nstructure:\n  - block: 01-b\n", encoding="utf-8")
    (lib / "blocks" / "01-b" / "block.yaml").write_text(
        "title: B\nduration_minutes: 45\nresources: []\n", encoding="utf-8"
    )
    (lib / "blocks" / "01-b" / "slides.md").write_text(
        "---\nsession: '01'\ntitle: T\nversion: 1.0.0\n---\n\n"
        "--- slide\nid: s-01\nlayout: Full\ntitle: T\n"
        "full:\n  type: ul\n  items:\n    - First\n---\n",
        encoding="utf-8",
    )
    out = tmp_path / "build"
    # Through the CLI, which is what derives the template and layout
    # map from the library rather than being handed them.
    corpus_main([str(lib), "--out", str(out)])
    deck = next((out / "sessions").rglob("*.pptx"))
    slide = Presentation(str(deck)).slides[0]
    for one in (_bullet_props(p) for p in _body_paragraphs(slide)):
        assert one == [], f"the template's own bullet was overridden: {one}"
